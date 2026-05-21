"""
生成周工作汇报 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import csv
import math
import os
from collections import defaultdict
from pathlib import Path

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案
COLOR_PRIMARY = RGBColor(0x1F, 0x49, 0x7D)   # 深蓝
COLOR_ACCENT = RGBColor(0x2E, 0x86, 0xAB)    # 浅蓝
COLOR_SUCCESS = RGBColor(0x27, 0xAE, 0x60)   # 绿色
COLOR_WARNING = RGBColor(0xE6, 0x7E, 0x22)   # 橙色
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)      # 深灰
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF8, 0xFA)  # 浅灰背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

STAGE1_PREDICTIONS_PATH = Path("outputs/stage1/latest_metrics/test_stage1_predictions.csv")


def _metric_summary(rows):
    """计算回归误差摘要，字段与 stage1 预测明细保持一致。"""
    if not rows:
        return {
            "n": 0,
            "mae": 0.0,
            "rmse": 0.0,
            "min_ae": 0.0,
            "max_ae": 0.0,
            "bias": 0.0,
            "err_min": 0.0,
            "err_max": 0.0,
            "within1": 0.0,
            "pred_min": 0.0,
            "pred_max": 0.0,
        }
    ae = [row["absolute_error"] for row in rows]
    se = [row["signed_error"] for row in rows]
    pred = [row["predicted_distance"] for row in rows]
    n = len(rows)
    return {
        "n": n,
        "mae": sum(ae) / n,
        "rmse": math.sqrt(sum(error * error for error in se) / n),
        "min_ae": min(ae),
        "max_ae": max(ae),
        "bias": sum(se) / n,
        "err_min": min(se),
        "err_max": max(se),
        "within1": sum(error <= 1.0 for error in ae) / n,
        "pred_min": min(pred),
        "pred_max": max(pred),
    }


def load_stage1_report_data(path=STAGE1_PREDICTIONS_PATH):
    """读取最新 stage1 预测明细，生成 PPT 所需的总体、分距离和文件级统计。"""
    if not path.exists():
        return None

    rows = []
    with path.open(newline="") as file:
        for row in csv.DictReader(file):
            parsed = dict(row)
            for key in ("target_distance", "predicted_distance", "signed_error", "absolute_error"):
                parsed[key] = float(parsed[key])
            rows.append(parsed)

    by_distance = defaultdict(list)
    by_file = defaultdict(list)
    for row in rows:
        by_distance[row["target_distance"]].append(row)
        by_file[row["raw_id"]].append(row)

    distance_rows = []
    for distance in sorted(by_distance):
        summary = _metric_summary(by_distance[distance])
        summary["distance"] = distance
        distance_rows.append(summary)

    file_rows = []
    for raw_id, file_segments in sorted(by_file.items()):
        target = file_segments[0]["target_distance"]
        prediction = sum(row["predicted_distance"] for row in file_segments) / len(file_segments)
        signed_error = prediction - target
        file_rows.append(
            {
                "raw_id": raw_id,
                "target_distance": target,
                "predicted_distance": prediction,
                "signed_error": signed_error,
                "absolute_error": abs(signed_error),
            }
        )

    return {
        "source_path": str(path),
        "segment": _metric_summary(rows),
        "file": _metric_summary(file_rows),
        "by_distance": distance_rows,
        "worst_segments": sorted(rows, key=lambda row: row["absolute_error"], reverse=True)[:6],
        "best_segments": sorted(rows, key=lambda row: row["absolute_error"])[:3],
        "worst_files": sorted(file_rows, key=lambda row: row["absolute_error"], reverse=True)[:5],
    }


STAGE1_REPORT = load_stage1_report_data()


def metric_value(key, level="segment", default=0.0):
    if STAGE1_REPORT is None:
        fallback = {
            ("segment", "mae"): 0.3703,
            ("segment", "rmse"): 0.5934,
            ("segment", "within1"): 0.9286,
            ("segment", "min_ae"): 0.0060,
            ("segment", "max_ae"): 2.4668,
            ("segment", "bias"): -0.1504,
            ("segment", "err_min"): -2.4668,
            ("segment", "err_max"): 1.1769,
            ("segment", "pred_min"): 1.9274,
            ("segment", "pred_max"): 16.1776,
            ("file", "mae"): 0.2534,
            ("file", "rmse"): 0.4603,
            ("file", "within1"): 0.9524,
            ("file", "max_ae"): 1.8852,
        }
        return fallback.get((level, key), default)
    return STAGE1_REPORT[level][key]

# ============================================================
# 辅助函数
# ============================================================

def add_slide_background(slide, color=COLOR_WHITE):
    """设置幻灯片背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header_bar(slide, title, subtitle=""):
    """添加顶部标题栏"""
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

def add_bullet_points(slide, points, left=0.5, top=1.6, width=12.333, font_size=18, line_spacing=1.3):
    """添加要点列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if point == "":
            p.text = ""
            p.space_after = Pt(6)
        elif point.startswith("  "):
            p.text = "    • " + point.strip()
            p.font.size = Pt(font_size - 2)
            p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        else:
            p.text = "• " + point
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)

def add_two_columns(slide, left_points, right_points, left_title="", right_title="", top=1.6):
    """添加双栏布局"""
    if left_title:
        txBox_title = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(5.8), Inches(0.4))
        tf = txBox_title.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        top_l = top + 0.5
    else:
        top_l = top

    txBox_l = slide.shapes.add_textbox(Inches(0.5), Inches(top_l), Inches(5.8), Inches(5))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    for i, point in enumerate(left_points):
        if i == 0:
            p = tf_l.paragraphs[0]
        else:
            p = tf_l.add_paragraph()
        p.text = "• " + point if point else ""
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)

    if right_title:
        txBox_title2 = slide.shapes.add_textbox(Inches(6.8), Inches(top), Inches(5.8), Inches(0.4))
        tf2 = txBox_title2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = right_title
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_PRIMARY
        top_r = top + 0.5
    else:
        top_r = top

    txBox_r = slide.shapes.add_textbox(Inches(6.8), Inches(top_r), Inches(5.8), Inches(5))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    for i, point in enumerate(right_points):
        if i == 0:
            p = tf_r.paragraphs[0]
        else:
            p = tf_r.add_paragraph()
        p.text = "• " + point if point else ""
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)

def add_card(slide, title, content, left, top, width=5.8, height=1.5, color=COLOR_LIGHT_BG):
    """添加卡片"""
    card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = COLOR_ACCENT

    title_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(width - 0.3), Inches(0.4))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    content_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.5), Inches(width - 0.3), Inches(height - 0.6))
    tf = content_box.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.text = content
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_TEXT

def add_section_divider(prs, num, title, subtitle=""):
    """添加章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_slide_background(slide)

    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # 装饰线
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(3.5), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"PART {num}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)
    p.font.bold = False

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE

    if subtitle:
        p3 = tf.add_paragraph()
        p3.text = subtitle
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)

    return slide

# ============================================================
# 封面
# ============================================================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)

# 背景大色块
shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(4.5))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_PRIMARY
shape.line.fill.background()

# 标题
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.5), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "本周工作汇报"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE

# 副标题
txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(11.5), Inches(0.6))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "基于双通道传感信号的管道泄漏检测与定位"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)

# 分隔线
line = slide.shapes.add_shape(1, Inches(0.7), Inches(3.4), Inches(2), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_ACCENT
line.line.fill.background()

# 元信息
txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11.5), Inches(1.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "汇报周期：2026-04-14 — 2026-04-22"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

p4 = tf3.add_paragraph()
p4.text = "姓名：[你的名字]    导师：[导师姓名]"
p4.font.size = Pt(16)
p4.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# ============================================================
# PART 1: 本周工作概览
# ============================================================
add_section_divider(prs, "01", "本周工作概览", "已完成的主要工作内容")

# 日程表页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "本周工作日程", "2026-04-14 — 2026-04-22")

items = [
    ("04-14", "新增 Conformer 模型，完成模型结构扩展", COLOR_ACCENT),
    ("04-16", "完成论文模板搭建，更新论文内容框架", COLOR_SUCCESS),
    ("04-17", "完善 .gitignore，优化项目工程化管理", COLOR_ACCENT),
    ("04-19", "新增 Stage2 混淆矩阵与 Macro-F1 评测逻辑", COLOR_SUCCESS),
    ("04-19", "修改论文正文，补充数据表格内容", COLOR_SUCCESS),
    ("04-19", "构建 LaTeX 编译命令与产物输出（Justfile）", COLOR_ACCENT),
    ("04-22", "补充完整实验结果表格，更新论文图表", COLOR_SUCCESS),
    ("04-22", "生成毕业设计答辩 PPT", COLOR_SUCCESS),
]

y = 1.55
for date, desc, color in items:
    # 日期标签
    tag = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(1.0), Inches(0.5))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()

    tag_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.08), Inches(1.0), Inches(0.4))
    p = tag_text.text_frame.paragraphs[0]
    p.text = date
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 描述
    desc_text = slide.shapes.add_textbox(Inches(1.7), Inches(y + 0.08), Inches(11), Inches(0.4))
    p2 = desc_text.text_frame.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_TEXT

    y += 0.72

# ============================================================
# PART 2: 模型设置
# ============================================================
add_section_divider(prs, "02", "模型设置", "两阶段网络架构与参数配置")

# 模型架构页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "模型架构 — 一维卷积特征提取网络")

# 左：架构说明
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(5.5), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "网络结构"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

points = [
    "每级模块：Conv + BN + GELU",
    "共 4 级卷积，通道逐层翻倍",
    "全局自适应平均池化",
    "全连接任务头",
    "",
    "Backbone 配置",
    "  32 → 64 → 128 → 256 通道",
    "  卷积核：15 / 9 / 7 / 5",
    "  步长：均为 4",
    "",
    "模型参数量",
    "  Stage2 分类：~28 万参数",
    "  Stage1 回归：~28 万参数",
]
for pt in points:
    p = tf.add_paragraph()
    if pt == "":
        p.text = ""
        p.space_after = Pt(4)
    elif pt.startswith("  "):
        p.text = "  " + pt.strip()
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    else:
        p.text = "• " + pt
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(6)

# 右：配置表
table = slide.shapes.add_table(5, 4, Inches(6.5), Inches(1.5), Inches(6.3), Inches(2.8)).table
headers = ["层级", "通道数", "卷积核", "步长"]
data = [
    ["Conv1", "32", "15", "4"],
    ["Conv2", "64", "9", "4"],
    ["Conv3", "128", "7", "4"],
    ["Conv4", "256", "5", "4"],
]

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER

# 两阶段任务卡片
add_card(slide, "Stage2 分类任务", "输入: 1×128000 单通道\n输出: 三分类标签\n损失: CrossEntropyLoss",
         left=6.5, top=4.5, width=3.0, height=1.3)
add_card(slide, "Stage1 回归任务", "输入: 2×128000 双通道\n输出: 泄漏距离（实数）\n损失: SmoothL1Loss",
         left=9.7, top=4.5, width=3.0, height=1.3)

# 两阶段对比
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "两阶段任务对比", "Stage2 分类 vs Stage1 回归")

add_two_columns(
    slide,
    left_points=[
        "输入: 单通道分段信号 (1×128000)",
        "任务: 三分类（正常+两类泄漏形态）",
        "损失函数: CrossEntropyLoss",
        "评测: Accuracy / Macro-F1 / 混淆矩阵",
        "片段级: 逐片段独立预测",
        "文件级: 同一文件片段 logits 聚合",
        "测试集 Acc: 0.9968",
        "测试集 Macro-F1: 0.9972",
        "文件级 Acc/F1: 1.0000",
    ],
    right_points=[
        "输入: 双通道信号对 (2×128000)",
        "任务: 泄漏距离回归",
        "损失函数: SmoothL1Loss",
        "评测: MAE / RMSE",
        "双通道保留泄漏传播位置差异",
        "逐样本标准化归一化",
        "验证集 MAE: 0.3148",
        "测试集 MAE: 0.3703",
        "测试集 RMSE: 0.5934",
    ],
    left_title="Stage2 分类",
    right_title="Stage1 回归",
    top=1.5
)

# ============================================================
# PART 3: 超参数
# ============================================================
add_section_divider(prs, "03", "超参数调节", "训练配置与调参决策")

slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "超参数配置", "Stage2 分类 / Stage1 回归")

# 超参数对比表
table = slide.shapes.add_table(9, 3, Inches(1), Inches(1.5), Inches(11), Inches(4.5)).table
headers = ["超参数", "Stage2 分类", "Stage1 回归"]
data = [
    ["初始学习率", "0.001", "0.001"],
    ["权重衰减", "0.0001", "0.0001"],
    ["Batch Size", "64", "32"],
    ["最大 Epochs", "40", "40"],
    ["早停 Patience", "8", "8"],
    ["梯度裁剪阈值", "1.0", "1.0"],
    ["学习率调度", "CosineAnnealing", "CosineAnnealing"],
    ["Dropout", "0.2", "0.2"],
]

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER

# 关键决策说明
txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "关键设计决策："
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

decisions = [
    "① Smooth L1 > MSE：误差 > 1 时线性惩罚，对离群点更鲁棒",
    "② Cosine Annealing：学习率从 0.001 平滑衰减至 1e-6，稳定收敛",
    "③ 梯度裁剪（1.0）：防止长序列训练中的梯度爆炸",
    "④ 早停（patience=8）：验证集指标未改善则提前停止，避免过拟合",
]
for d in decisions:
    p = tf.add_paragraph()
    p.text = d
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(6)

# ============================================================
# PART 4: 论文内容
# ============================================================
add_section_divider(prs, "04", "论文内容更新", "各章节完成情况与新增内容")

slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "论文各章节完成情况", "截至 2026-04-22")

chapters = [
    ("第一章 绪论", "✅ 完成", COLOR_SUCCESS, [
        "研究背景：管道安全重要性、泄漏故障特点",
        "研究问题：双通道时序信号 → 状态识别+距离估计",
        "本文工作：数据构建、任务拆分、模型设计、训练框架",
    ]),
    ("第二章 数据集构建", "✅ 完成", COLOR_SUCCESS, [
        "原始数据组织：CSV 分目录（leak / normal）",
        "标签映射：Stage2 分类 + Stage1 回归",
        "分段策略：25600 Hz × 5 s = 128000 点",
        "数据集统计：Stage2 3024 样本，Stage1 1122 样本",
    ]),
    ("第三章 模型设计", "✅ 完成", COLOR_SUCCESS, [
        "总体框架：数据切分 → 分类识别 → 距离估计",
        "一维卷积网络：4 级 Conv + BN + GELU",
        "分类头（交叉熵）+ 回归头（Smooth L1）",
        "训练策略：AdamW + Cosine + 早停 + 梯度裁剪",
    ]),
    ("第四章 实验结果", "✅ 完成", COLOR_SUCCESS, [
        "双层评测：片段级 + 文件级聚合",
        "Stage2 Acc=0.9968，文件级 Acc=1.0",
        "Stage1 MAE=0.3703，RMSE=0.5934",
        "混淆矩阵分析：仅 1 次片段级误分类",
    ]),
    ("第五章 结论展望", "✅ 完成", COLOR_SUCCESS, [
        "总结两阶段框架有效性",
        "展望：数据扩充、模型改进、联合学习、工程部署",
    ]),
]

y = 1.45
for title, status, color, points in chapters:
    # 章节标题行
    tag = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(2.2), Inches(0.42))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()

    tag_text = slide.shapes.add_textbox(Inches(0.45), Inches(y + 0.07), Inches(2.1), Inches(0.35))
    p = tag_text.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    status_text = slide.shapes.add_textbox(Inches(2.8), Inches(y + 0.07), Inches(0.8), Inches(0.35))
    p2 = status_text.text_frame.paragraphs[0]
    p2.text = status
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = color

    # 内容
    content = "  |  ".join(points)
    content_text = slide.shapes.add_textbox(Inches(3.8), Inches(y + 0.05), Inches(9.2), Inches(0.38))
    p3 = content_text.text_frame.paragraphs[0]
    p3.text = content
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    y += 0.58

# 新增内容页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "本周新增论文内容", "重点更新模块")

new_contents = [
    ("文件级聚合评测", "同一原始文件的多个片段 logits 取平均后再投票，有效避免片段级虚高，验证真正的泛化能力", COLOR_ACCENT),
    ("Macro-F1 指标", "各类别 F1 等权平均，应对 Stage2 分类任务中类别不均衡（780/1440/804）问题", COLOR_ACCENT),
    ("Smooth L1 损失", "Stage1 回归任务采用，对大误差（|e|≥1）采用线性惩罚，对离群预测更鲁棒", COLOR_ACCENT),
    ("LaTeX 编译流程", "Justfile 中新增 thesis-build / thesis-clean / thesis-rebuild 命令，自动化论文编译", COLOR_SUCCESS),
]

y = 1.5
for title, desc, color in new_contents:
    card = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.333), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = color

    # 左侧色条
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(0.12), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.12), Inches(11.5), Inches(0.4))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.55), Inches(11.5), Inches(0.6))
    p2 = desc_box.text_frame.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT

    y += 1.35

# ============================================================
# PART 5: 实验环节
# ============================================================
add_section_divider(prs, "05", "实验环节", "评测体系设计与结果分析")

# 评测体系页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "评测体系设计", "片段级 + 文件级双层评测")

add_two_columns(
    slide,
    left_points=[
        "片段级评测",
        "",
        "逐片段独立预测",
        "直接计算 Accuracy / Macro-F1",
        "每个片段是独立样本",
        "",
        "指标：",
        "  Accuracy = N_correct / N_total",
        "  Macro-F1 = 等权平均各类别 F1",
        "  混淆矩阵：展示类别混淆关系",
    ],
    right_points=[
        "文件级评测（新增）",
        "",
        "同一文件所有片段 logits 聚合",
        "取平均后 argmax 投票",
        "更严格，考验跨片段泛化能力",
        "",
        "目的：",
        "  避免同一文件片段相似性造成的",
        "  片段级虚高",
        "  验证模型真正见过未见文件的能力",
    ],
    left_title="片段级",
    right_title="文件级（新增）",
    top=1.5
)

# 实验结果页 - Stage2
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "实验结果 — Stage2 分类任务", "片段级 + 文件级双层结果")

# 结果卡片
results_s2 = [
    ("片段级 Accuracy", "0.9968", "312 测试样本", COLOR_ACCENT),
    ("片段级 Macro-F1", "0.9972", "仅 1 次误分类", COLOR_ACCENT),
    ("文件级 Accuracy", "1.0000", "26 个文件全部正确", COLOR_SUCCESS),
    ("文件级 Macro-F1", "1.0000", "主对角线全正", COLOR_SUCCESS),
]

y = 1.5
for i, (title, value, note, color) in enumerate(results_s2):
    left = 0.5 + (i % 2) * 6.4
    top = 1.5 + (i // 2) * 1.6

    card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(6.0), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(5.5), Inches(0.4))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    value_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.55), Inches(5.5), Inches(0.5))
    p2 = value_box.text_frame.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE

    note_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.0), Inches(5.5), Inches(0.35))
    p3 = note_box.text_frame.paragraphs[0]
    p3.text = note
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)

# 混淆矩阵
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(0.4))
p = txBox.text_frame.paragraphs[0]
p.text = "混淆矩阵（片段级）："
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

table = slide.shapes.add_table(4, 4, Inches(3.5), Inches(5.1), Inches(6.5), Inches(2.0)).table
cm_headers = ["真实\\预测", "0", "1", "2"]
cm_data = [
    ["0", "60", "0", "0"],
    ["1", "0", "155", "1"],
    ["2", "0", "0", "96"],
]

for i, h in enumerate(cm_headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

for row_idx, row_data in enumerate(cm_data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER
        if col_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xEE, 0xF3, 0xF7)

# 实验结果页 - Stage1
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "实验结果 — Stage1 回归任务", "双通道距离估计")

# MAE 对比
card_mae = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.2))
card_mae.fill.solid()
card_mae.fill.fore_color.rgb = COLOR_PRIMARY
card_mae.line.fill.background()

txBox1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.65), Inches(5.4), Inches(0.5))
p = txBox1.text_frame.paragraphs[0]
p.text = "测试集 MAE"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)

txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(5.4), Inches(0.8))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = f"{metric_value('mae'):.4f}"
p2.font.size = Pt(48)
p2.font.bold = True
p2.font.color.rgb = COLOR_WHITE

txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(2.95), Inches(5.4), Inches(0.5))
p3 = txBox3.text_frame.paragraphs[0]
p3.text = f"均方根误差 RMSE: {metric_value('rmse'):.4f}"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)

# 基线对比
card_base = slide.shapes.add_shape(1, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.2))
card_base.fill.solid()
card_base.fill.fore_color.rgb = RGBColor(0xEE, 0x7E, 0x22)
card_base.line.fill.background()

txBox4 = slide.shapes.add_textbox(Inches(7.0), Inches(1.65), Inches(5.4), Inches(0.5))
p4 = txBox4.text_frame.paragraphs[0]
p4.text = "vs 统计基线"
p4.font.size = Pt(18)
p4.font.color.rgb = COLOR_WHITE

txBox5 = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.4), Inches(0.8))
p5 = txBox5.text_frame.paragraphs[0]
p5.text = "MAE > 3"
p5.font.size = Pt(48)
p5.font.bold = True
p5.font.color.rgb = COLOR_WHITE

txBox6 = slide.shapes.add_textbox(Inches(7.0), Inches(2.95), Inches(5.4), Inches(0.5))
p6 = txBox6.text_frame.paragraphs[0]
p6.text = "均值/中位数预测 基线对比"
p6.font.size = Pt(14)
p6.font.color.rgb = RGBColor(0xFF, 0xDD, 0xAA)

# 说明
txBox7 = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(3))
tf = txBox7.text_frame
tf.word_wrap = True
p7 = tf.paragraphs[0]
p7.text = "结果分析："
p7.font.size = Pt(16)
p7.font.bold = True
p7.font.color.rgb = COLOR_PRIMARY

points = [
    f"测试集 MAE 仅为 {metric_value('mae'):.4f}，相比基线（>3）有数量级优势",
    f"误差 ≤ 1 个距离单位的片段占比为 {metric_value('within1') * 100:.1f}%",
    f"文件级聚合后 MAE 降至 {metric_value('mae', 'file'):.4f}，说明多片段平均能进一步稳定定位",
    f"最大绝对误差为 {metric_value('max_ae'):.4f}，主要来自个别距离档位的低估",
    "说明模型有效捕获了双通道信号模式与泄漏距离的映射关系",
]
for pt in points:
    p = tf.add_paragraph()
    p.text = "• " + pt
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(10)

# Stage1 物理误差分析页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "Stage1 物理误差分析", "最大/最小误差 · 误差方向 · 文件级聚合")

metric_cards = [
    ("最小绝对误差", f"{metric_value('min_ae'):.4f}", "单个片段几乎完全命中"),
    ("最大绝对误差", f"{metric_value('max_ae'):.4f}", "用于定位最坏预测案例"),
    ("最大低估", f"{metric_value('err_min'):.4f}", "prediction - target"),
    ("最大高估", f"{metric_value('err_max'):.4f}", "prediction - target"),
    ("预测范围", f"{metric_value('pred_min'):.2f} ~ {metric_value('pred_max'):.2f}", "真实距离范围 2 ~ 16"),
    ("文件级 MAE", f"{metric_value('mae', 'file'):.4f}", "同一 ABC 多片段取均值"),
]

for idx, (title, value, note) in enumerate(metric_cards):
    col = idx % 3
    row_idx = idx // 3
    left = 0.55 + col * 4.25
    top = 1.55 + row_idx * 1.65
    card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(3.75), Inches(1.25))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = COLOR_ACCENT

    title_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(3.4), Inches(0.28))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    value_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.44), Inches(3.4), Inches(0.45))
    p2 = value_box.text_frame.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(24 if len(value) < 10 else 18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WARNING if "低估" in title or "最大" in title else COLOR_SUCCESS

    note_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.9), Inches(3.4), Inches(0.25))
    p3 = note_box.text_frame.paragraphs[0]
    p3.text = note
    p3.font.size = Pt(10)
    p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

summary_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.05), Inches(12), Inches(1.25))
tf = summary_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "汇报解读："
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY
interpretations = [
    f"片段级 {metric_value('within1') * 100:.1f}% 的预测误差不超过 1 个距离单位，说明大多数定位结果具备工程可读性",
    f"文件级聚合后 MAE 从 {metric_value('mae'):.4f} 降至 {metric_value('mae', 'file'):.4f}，多片段平均能抑制单段波动",
    "最大低估比最大高估更明显，后续应重点分析高距离档位信号特征是否不足",
]
for item in interpretations:
    p = tf.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(4)

# Stage1 分距离与异常样本页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "Stage1 误差定位", "按真实距离档位拆解 + 最大误差样本")

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(6.2), Inches(0.35))
p = txBox.text_frame.paragraphs[0]
p.text = "按距离档位误差"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

distance_data = STAGE1_REPORT["by_distance"] if STAGE1_REPORT else [
    {"distance": 2, "n": 6, "mae": 0.0750, "max_ae": 0.1482, "bias": 0.0277, "within1": 1.0},
    {"distance": 3, "n": 24, "mae": 0.2519, "max_ae": 0.8743, "bias": 0.0756, "within1": 1.0},
    {"distance": 12, "n": 30, "mae": 0.3519, "max_ae": 2.4668, "bias": -0.1672, "within1": 0.9667},
    {"distance": 16, "n": 12, "mae": 1.1186, "max_ae": 2.2493, "bias": -1.0627, "within1": 0.5},
]

table = slide.shapes.add_table(len(distance_data) + 1, 5, Inches(0.5), Inches(1.9), Inches(6.25), Inches(3.8)).table
headers = ["距离", "样本", "MAE", "最大误差", "Bias"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

for row_idx, row_data in enumerate(distance_data):
    values = [
        f"{row_data['distance']:.0f}",
        f"{row_data['n']:.0f}",
        f"{row_data['mae']:.3f}",
        f"{row_data['max_ae']:.3f}",
        f"{row_data['bias']:.3f}",
    ]
    for col_idx, value in enumerate(values):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = value
        if row_data["mae"] > 1.0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xEE, 0xDD)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(7.05), Inches(1.45), Inches(5.8), Inches(0.35))
p = txBox2.text_frame.paragraphs[0]
p.text = "最大误差样本"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

worst_segments = STAGE1_REPORT["worst_segments"] if STAGE1_REPORT else []
worst_table = slide.shapes.add_table(7, 5, Inches(7.05), Inches(1.9), Inches(5.75), Inches(3.8)).table
headers = ["文件", "真实", "预测", "误差", "|误差|"]
for col_idx, header in enumerate(headers):
    cell = worst_table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_WARNING
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

for row_idx in range(6):
    row_data = worst_segments[row_idx] if row_idx < len(worst_segments) else {
        "raw_id": "",
        "target_distance": 0,
        "predicted_distance": 0,
        "signed_error": 0,
        "absolute_error": 0,
    }
    values = [
        row_data["raw_id"],
        f"{row_data['target_distance']:.0f}" if row_data["raw_id"] else "",
        f"{row_data['predicted_distance']:.2f}" if row_data["raw_id"] else "",
        f"{row_data['signed_error']:.2f}" if row_data["raw_id"] else "",
        f"{row_data['absolute_error']:.2f}" if row_data["raw_id"] else "",
    ]
    for col_idx, value in enumerate(values):
        cell = worst_table.cell(row_idx + 1, col_idx)
        cell.text = value
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER

note_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.65))
p = note_box.text_frame.paragraphs[0]
p.text = "关键发现：16 距离档位存在系统性低估；最大单段误差来自 ABC155（真实 12，预测 9.53），最大文件级误差来自 ABC162（真实 16，均值预测 14.11）。"
p.font.size = Pt(13)
p.font.color.rgb = COLOR_TEXT

# 结果分析页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "实验结论与分析", "两阶段任务综合分析")

add_two_columns(
    slide,
    left_points=[
        "Stage2 分类 — 结论",
        "",
        "一维卷积能有效捕获泄漏形态差异",
        "文件级 1.0 排除片段级虚高结论",
        "26 个测试文件全部正确分类",
        "仅 1 次片段级误分类（类别1→2）",
        "Macro-F1=0.9972 表明各类别均稳定",
        "",
        "局限性",
        "随机划分测试集规模有限",
        "工况相似性可能影响泛化结论",
    ],
    right_points=[
        "Stage1 回归 — 结论",
        "",
        "双通道信号提供有效位置信息",
        "误差控制在相邻距离档位",
        "MAE 0.37 远优于统计基线 >3",
        "RMSE 0.59 对较大误差敏感但仍优",
        "",
        "后续改进方向",
        "更严格的数据切分策略",
        "补充跨工况泛化实验",
        "不同网络结构对比",
    ],
    left_title="Stage2 分类",
    right_title="Stage1 回归",
    top=1.5
)

# ============================================================
# PART 6: 接下来工作
# ============================================================
add_section_divider(prs, "06", "接下来工作计划", "近期 / 中期 / 长期目标")

slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "工作计划 — 三阶段目标", "近期 · 中期 · 长期")

# 三个阶段卡片
phases = [
    ("近期目标（1-2 周）", COLOR_ACCENT, [
        ("跨工况泛化实验", "高优先级", "设计更严格的数据切分策略"),
        ("模型对比实验", "中优先级", "尝试 ResNet1D / TCN 等 backbone"),
        ("章节内容细化", "中优先级", "补充消融实验、误差分析小节"),
    ]),
    ("中期目标（1 个月）", COLOR_PRIMARY, [
        ("模型结构改进", "中优先级", "残差连接 / 注意力机制"),
        ("时频联合建模", "低优先级", "STFT + 2D CNN 探索"),
        ("联合学习框架", "低优先级", "分类+回归联合优化"),
    ]),
    ("长期目标（毕业前）", COLOR_SUCCESS, [
        ("论文定稿", "必做", "补充实验、润色语言、查重"),
        ("工程部署", "选做", "模型轻量化、实时推理"),
        ("答辩准备", "必做", "预演、问答准备"),
    ]),
]

y = 1.5
for phase_title, color, items in phases:
    # 阶段标题
    bar = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.1), Inches(12), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    p.text = phase_title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    y += 0.6

    for item_title, priority, desc in items:
        # 优先级标签
        tag_color = COLOR_ACCENT if priority == "高优先级" else (RGBColor(0xAA, 0xAA, 0xAA) if priority == "低优先级" else RGBColor(0x77, 0x88, 0x99))
        tag = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(1.2), Inches(0.38))
        tag.fill.solid()
        tag.fill.fore_color.rgb = tag_color
        tag.line.fill.background()

        tag_t = slide.shapes.add_textbox(Inches(0.55), Inches(y + 0.07), Inches(1.1), Inches(0.3))
        p = tag_t.text_frame.paragraphs[0]
        p.text = priority
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

        # 任务标题
        title_t = slide.shapes.add_textbox(Inches(1.9), Inches(y + 0.05), Inches(3.5), Inches(0.35))
        p2 = title_t.text_frame.paragraphs[0]
        p2.text = item_title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT

        # 描述
        desc_t = slide.shapes.add_textbox(Inches(5.5), Inches(y + 0.05), Inches(7.5), Inches(0.35))
        p3 = desc_t.text_frame.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        y += 0.48

    y += 0.15

# 问题与风险页
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)
add_header_bar(slide, "当前问题与风险", "需要关注的事项")

risks = [
    ("测试集规模有限", "泛化结论说服力不足", "设计跨文件/跨工况划分实验，补充更多测试数据", "🔴 高"),
    ("类别不均衡", "少数类（类别0）识别可能偏弱", "尝试加权交叉熵或过采样策略", "🟡 中"),
    ("论文篇幅", "规范要求与内容深度需平衡", "按模板精炼语言，图表互补表达", "🟡 中"),
    ("时间进度", "毕业节点临近", "按近期→中期→长期优先级推进", "🔴 高"),
]

y = 1.5
for risk, impact, solution, level in risks:
    card = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.333), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = COLOR_ACCENT

    # 风险名称
    risk_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.12), Inches(3), Inches(0.4))
    p = risk_box.text_frame.paragraphs[0]
    p.text = risk
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # 影响
    impact_box = slide.shapes.add_textbox(Inches(3.8), Inches(y + 0.12), Inches(3.5), Inches(0.4))
    p2 = impact_box.text_frame.paragraphs[0]
    p2.text = "影响：" + impact
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(0x88, 0x44, 0x44)

    # 应对
    solution_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.55), Inches(11.5), Inches(0.5))
    p3 = solution_box.text_frame.paragraphs[0]
    p3.text = "应对：" + solution
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT

    # 等级
    level_box = slide.shapes.add_textbox(Inches(11.5), Inches(y + 0.12), Inches(1.5), Inches(0.4))
    p4 = level_box.text_frame.paragraphs[0]
    p4.text = level
    p4.font.size = Pt(13)
    p4.font.bold = True

    y += 1.35

# ============================================================
# 感谢页
# ============================================================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)
add_slide_background(slide)

shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_PRIMARY
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0), Inches(2.8), Inches(13.333), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "谢谢！"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "欢迎提问与交流"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)
p2.alignment = PP_ALIGN.CENTER

line = slide.shapes.add_shape(1, Inches(5.5), Inches(4.0), Inches(2.3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = COLOR_ACCENT
line.line.fill.background()

# 保存
output_path = "outputs/本周工作汇报PPT.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
