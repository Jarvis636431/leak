"""Generate a placeholder presentation for the Internoise experiment refresh."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("outputs/本周工作汇报PPT_placeholder.pptx")


def add_centered_text(slide, title: str, subtitle: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.2))
    frame = box.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = frame.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(17)
    p2.alignment = PP_ALIGN.CENTER


def main() -> None:
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_text(
        slide,
        "Internoise 实验汇报占位",
        "旧实验指标已清理；Stage2/Stage1 新训练完成后重新生成正式 PPT。",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_centered_text(
        slide,
        "待补充内容",
        "数据集统计、Stage2 五分类结果、Stage1 回归结果、混淆矩阵和训练曲线。",
    )

    prs.save(OUTPUT_PATH)
    print(f"Placeholder PPT written to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
