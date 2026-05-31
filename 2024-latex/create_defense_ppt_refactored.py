#!/usr/bin/env python3
"""Generate a refactored graduation defense deck for the leak detection project."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
HERE = Path(__file__).resolve().parent
OUT = HERE / "defense_presentation_refactored.pptx"

W = Inches(13.333)
H = Inches(7.5)

COLORS = {
    "ink": RGBColor(24, 31, 45),
    "navy": RGBColor(0, 74, 142),
    "blue": RGBColor(0, 91, 172),
    "teal": RGBColor(0, 133, 194),
    "amber": RGBColor(60, 160, 220),
    "red": RGBColor(200, 70, 62),
    "paper": RGBColor(247, 251, 255),
    "line": RGBColor(214, 228, 242),
    "muted": RGBColor(93, 105, 122),
    "white": RGBColor(255, 255, 255),
}

FONT = "Microsoft YaHei"
TOTAL = 22


def cm(v: float) -> int:
    return Inches(v / 2.54)


def display_size(size: int | float) -> int | float:
    """Slightly enlarge generated text for projector readability."""
    if size <= 8:
        return size + 1
    if size <= 13:
        return size + 2
    if size <= 18:
        return size + 2
    if size <= 25:
        return size + 2
    return size + 3


def add_rect(slide, x, y, w, h, color, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    return s


def add_text(slide, x, y, w, h, text, size=16, color=None, bold=False, align=PP_ALIGN.LEFT):
    rendered_size = display_size(size)
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(rendered_size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    p.alignment = align
    p.space_after = Pt(0)
    return box


def add_lines(slide, x, y, w, h, lines, size=13, color=None, spacing=1.25, bold_first=False):
    rendered_size = display_size(size)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(rendered_size)
        p.font.color.rgb = color or COLORS["ink"]
        p.font.bold = bold_first and i == 0
        p.line_spacing = Pt(rendered_size * spacing)
        p.space_after = Pt(2)
    return box


def add_brand_corner(slide, dark=False):
    logo = HERE / ("天津大学_白色logo.png" if dark else "天津大学_logo2.png")
    add_image(slide, logo, Inches(10.15), Inches(0.32), Inches(2.35), Inches(0.48), False)


def add_title(slide, title, kicker=None, page=None):
    add_rect(slide, 0, 0, W, Inches(0.12), COLORS["amber"])
    add_brand_corner(slide)
    if kicker:
        add_text(slide, Inches(0.58), Inches(0.32), Inches(9.2), Inches(0.28), kicker, 10, COLORS["teal"], True)
    add_text(slide, Inches(0.55), Inches(0.55), Inches(9.45), Inches(0.55), title, 25, COLORS["ink"], True)
    add_rect(slide, Inches(0.58), Inches(1.18), Inches(1.35), Inches(0.035), COLORS["teal"])
    if page is not None:
        add_text(slide, Inches(12.25), Inches(7.05), Inches(0.7), Inches(0.2), f"{page}/{TOTAL}", 9, COLORS["muted"], False, PP_ALIGN.RIGHT)


def add_image(slide, path: Path, x, y, w, h, border=True):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    px = x + int((w - nw) / 2)
    py = y + int((h - nh) / 2)
    pic = slide.shapes.add_picture(str(path), px, py, width=nw, height=nh)
    return pic


def metric_card(slide, x, y, label, value, note, color):
    add_rect(slide, x, y, Inches(2.1), Inches(1.16), COLORS["white"], COLORS["line"], radius=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.12), Inches(1.8), Inches(0.2), label, 9, COLORS["muted"], True)
    add_text(slide, x + Inches(0.15), y + Inches(0.35), Inches(1.8), Inches(0.38), value, 22, color, True)
    add_text(slide, x + Inches(0.15), y + Inches(0.84), Inches(1.8), Inches(0.2), note, 8, COLORS["muted"])


def insight_box(slide, x, y, w, h, title, body, color=COLORS["navy"]):
    add_rect(slide, x, y, w, h, COLORS["white"], COLORS["line"], True)
    add_text(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.24), title, 12, color, True)
    add_text(slide, x + Inches(0.18), y + Inches(0.48), w - Inches(0.36), h - Inches(0.62), body, 10, COLORS["ink"])


def small_table(slide, x, y, widths, rows, header_color=COLORS["navy"], size=10):
    row_h = Inches(0.39)
    for r, row in enumerate(rows):
        cx = x
        for c, txt in enumerate(row):
            w = Inches(widths[c])
            color = header_color if r == 0 else (COLORS["paper"] if r % 2 else COLORS["white"])
            add_rect(slide, cx, y + r * row_h, w, row_h, color, COLORS["line"])
            add_text(
                slide,
                cx + Inches(0.03),
                y + r * row_h + Inches(0.045),
                w - Inches(0.06),
                row_h - Inches(0.06),
                txt,
                size,
                COLORS["white"] if r == 0 else COLORS["ink"],
                r == 0,
                PP_ALIGN.CENTER,
            )
            cx += w


def bar(slide, x, y, label, value, max_value, color, suffix=""):
    add_text(slide, x, y, Inches(1.2), Inches(0.25), label, 10, COLORS["ink"], True)
    add_rect(slide, x + Inches(1.25), y + Inches(0.06), Inches(3.25), Inches(0.13), COLORS["line"], None, True)
    add_rect(slide, x + Inches(1.25), y + Inches(0.06), Inches(3.25 * value / max_value), Inches(0.13), color, None, True)
    add_text(slide, x + Inches(4.62), y - Inches(0.005), Inches(0.8), Inches(0.22), f"{value:.4g}{suffix}", 9, COLORS["muted"], False, PP_ALIGN.RIGHT)


def add_footer(slide):
    add_rect(slide, 0, Inches(7.34), W, Inches(0.16), COLORS["navy"])


def make_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, COLORS["paper"])
    add_rect(s, 0, 0, W, Inches(0.22), COLORS["navy"])
    add_rect(s, 0, Inches(7.14), W, Inches(0.36), COLORS["navy"])
    add_rect(s, Inches(5.28), Inches(0.82), Inches(2.78), Inches(2.78), COLORS["white"], COLORS["line"], True)
    add_image(s, HERE / "figures/tjulogo.png", Inches(5.55), Inches(1.05), Inches(2.24), Inches(2.24), False)
    add_text(s, Inches(0.95), Inches(3.92), Inches(11.45), Inches(1.15), "基于双通道时序信号的\n两阶段管道泄漏检测与定位方法", 32, COLORS["navy"], True, PP_ALIGN.CENTER)
    add_rect(s, Inches(4.55), Inches(5.42), Inches(4.25), Inches(0.045), COLORS["blue"])
    add_text(s, Inches(1.45), Inches(5.73), Inches(10.45), Inches(0.35), "从“是否泄漏”到“泄漏位置”：分类先行，双通道回归定位", 17, COLORS["ink"], True, PP_ALIGN.CENTER)
    add_lines(s, Inches(3.55), Inches(6.35), Inches(6.25), Inches(0.52), ["天津大学 2026 届本科生毕业设计答辩", "答辩人：孟达  |  导师：顾鹏 讲师"], 11, COLORS["muted"], 1.35)
    add_footer(s)

    # 2
    s = prs.slides.add_slide(blank)
    add_title(s, "项目判断：这份 PPT 需要从说明书改成答辩故事", "RESTRUCTURE LOGIC", 2)
    add_lines(s, Inches(0.7), Inches(1.65), Inches(4.4), Inches(4.9), [
        "原稿主要问题",
        "1. 20 页几乎全是文字和手绘表格，图像证据没有进入主叙事",
        "2. 分类结果、定位误差、消融实验分散，听众难形成一个结论",
        "3. 方法页过于解释结构细节，缺少“为什么这样设计”的视觉证据",
        "4. 结论页同时讲贡献和局限，但没有回扣实验指标",
    ], 13, COLORS["ink"], 1.45, True)
    for i, (title, body, color) in enumerate([
        ("主线", "工程问题 → 两阶段任务 → 模型与数据 → 实验验证 → 边界与展望", COLORS["teal"]),
        ("证据", "每 1 页只保留 1 个判断，用图或指标支撑", COLORS["blue"]),
        ("节奏", "背景 3 页、方法 5 页、结果 5 页、总结 2 页", COLORS["amber"]),
    ]):
        y = Inches(1.65 + i * 1.35)
        add_rect(s, Inches(5.55), y, Inches(6.85), Inches(1.05), COLORS["white"], COLORS["line"], True)
        add_text(s, Inches(5.82), y + Inches(0.13), Inches(1.2), Inches(0.32), title, 15, color, True)
        add_text(s, Inches(7.0), y + Inches(0.16), Inches(5.05), Inches(0.55), body, 13, COLORS["ink"])
    add_footer(s)

    # 3
    s = prs.slides.add_slide(blank)
    add_title(s, "研究背景：泄漏检测的难点不是“有无信号”，而是噪声下可定位", "BACKGROUND", 3)
    add_image(s, ROOT / "管道泄露危害场景示意图.png", Inches(0.7), Inches(1.55), Inches(5.65), Inches(4.1))
    add_image(s, ROOT / "管道泄露声信号传播与双通道响应差异.png", Inches(6.65), Inches(1.55), Inches(5.95), Inches(4.1))
    add_lines(s, Inches(0.95), Inches(5.88), Inches(11.2), Inches(0.85), [
        "研究切入点：泄漏扰动会在左右传感器中形成幅值、相位和到达响应差异；分类负责过滤状态，双通道回归负责利用差异定位。"
    ], 15, COLORS["ink"], 1.2)
    add_footer(s)

    # 4
    s = prs.slides.add_slide(blank)
    add_title(s, "数据设计：先防止信息泄漏，再谈模型指标", "DATA PIPELINE", 4)
    add_image(s, ROOT / "实验台与数据采集流程图.png", Inches(0.65), Inches(1.45), Inches(5.75), Inches(3.0))
    add_image(s, ROOT / "数据采集到训练样本生成.png", Inches(6.75), Inches(1.45), Inches(5.85), Inches(3.0))
    metric_card(s, Inches(0.9), Inches(4.85), "Stage2 分类", "11,704", "单通道片段样本", COLORS["blue"])
    metric_card(s, Inches(3.25), Inches(4.85), "Stage1 回归", "3,737", "左右通道配对样本", COLORS["teal"])
    metric_card(s, Inches(5.6), Inches(4.85), "采样窗口", "5 s", "128,000 点/通道", COLORS["amber"])
    add_lines(s, Inches(8.05), Inches(4.72), Inches(4.35), Inches(1.45), [
        "关键处理原则",
        "同一原始文件的所有片段只进入同一个 train/val/test 集合，避免相邻片段跨集合导致测试指标虚高。",
    ], 12, COLORS["ink"], 1.35, True)
    add_footer(s)

    # 5
    s = prs.slides.add_slide(blank)
    add_title(s, "任务拆分：分类解决状态，回归只在泄漏样本上估计距离", "TWO-STAGE TASK", 5)
    add_image(s, ROOT / "stage1 与 stage2 标签映射示意.png", Inches(0.7), Inches(1.4), Inches(5.7), Inches(3.35))
    add_image(s, ROOT / "stage1 双通道距离回归机制.png", Inches(6.75), Inches(1.4), Inches(5.8), Inches(3.35))
    small_table(s, Inches(0.8), Inches(5.1), [1.5, 2.25, 2.0, 2.2, 2.2], [
        ["阶段", "输入", "输出", "损失函数", "主指标"],
        ["Stage2", "单通道 1 x 128000", "5 类状态", "Cross-Entropy", "Acc / Macro-F1"],
        ["Stage1", "双通道 2 x 128000", "距离 m", "Smooth L1", "MAE / RMSE"],
    ], size=9)
    add_footer(s)

    # 6
    s = prs.slides.add_slide(blank)
    add_title(s, "基础模型：1D CNN 用逐层降采样把长时序压到可学习表示", "MODEL BACKBONE", 6)
    add_image(s, ROOT / "一维 CNN 网络结构图.png", Inches(0.75), Inches(1.35), Inches(5.7), Inches(4.4))
    add_image(s, ROOT / "一维cnn 多尺度感受野示意图.png", Inches(6.85), Inches(1.35), Inches(5.55), Inches(2.9))
    add_lines(s, Inches(6.95), Inches(4.58), Inches(5.25), Inches(1.35), [
        "设计要点",
        "4 级卷积逐步降采样，先提取局部波形细节，再扩大感受野；全局平均池化减少参数量，任务头分别用于分类和回归。",
    ], 13, COLORS["ink"], 1.4, True)
    add_footer(s)

    # 7
    s = prs.slides.add_slide(blank)
    add_title(s, "结构对比：Conformer 增强长程建模，但不一定改善距离回归", "ARCHITECTURE COMPARISON", 7)
    add_image(s, ROOT / "conformer 时序建模结构图.png", Inches(0.7), Inches(1.35), Inches(5.8), Inches(3.2))
    add_image(s, ROOT / "conformer 编码块内部结构.png", Inches(6.75), Inches(1.35), Inches(5.85), Inches(3.2))
    add_lines(s, Inches(0.95), Inches(5.0), Inches(11.1), Inches(0.75), [
        "解释逻辑：Conformer 在 CNN 前端降采样后引入注意力，适合补足分类任务的全局上下文；但距离回归更依赖细粒度相位/幅值差异，强全局建模可能不稳定。"
    ], 14, COLORS["ink"], 1.25)
    add_footer(s)

    # 8
    s = prs.slides.add_slide(blank)
    add_title(s, "评价体系：同时看片段级、文件级和物理误差", "EVALUATION", 8)
    add_image(s, ROOT / "实验评价体系总览.png", Inches(0.75), Inches(1.35), Inches(6.1), Inches(4.15))
    add_image(s, ROOT / "文件级划分与片段级:文件级评测关系.png", Inches(7.15), Inches(1.35), Inches(5.25), Inches(4.15))
    add_lines(s, Inches(0.9), Inches(5.78), Inches(11.1), Inches(0.62), [
        "答辩中建议强调：片段级反映单窗口即时判断，文件级更接近连续采集决策；回归还必须报告 Bias、P95 和按距离档位误差。"
    ], 13, COLORS["ink"], 1.2)
    add_footer(s)

    # 9
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage2 结果：分类已接近工程可用，文件级聚合显著稳定", "CLASSIFICATION RESULT", 9)
    metric_card(s, Inches(0.75), Inches(1.5), "CNN 片段 Acc", "0.9415", "Macro-F1 0.9180", COLORS["blue"])
    metric_card(s, Inches(3.1), Inches(1.5), "CNN 文件 Acc", "0.9783", "Macro-F1 0.9836", COLORS["teal"])
    metric_card(s, Inches(5.45), Inches(1.5), "Conformer 片段 Acc", "0.9900", "Macro-F1 0.9883", COLORS["amber"])
    add_lines(s, Inches(8.0), Inches(1.5), Inches(4.25), Inches(0.95), [
        "文件级聚合后仅 1 个文件误判，说明连续采集场景下分类决策更稳定。"
    ], 12, COLORS["ink"], 1.25)
    add_image(s, HERE / "figures/experiments/stage2_segment_confusion_matrix.png", Inches(0.75), Inches(3.0), Inches(5.6), Inches(3.35))
    add_image(s, HERE / "figures/experiments/stage2_file_confusion_matrix.png", Inches(6.85), Inches(3.0), Inches(5.45), Inches(3.35))
    add_footer(s)

    # 10
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage2 对比：深度模型提升片段判别，传统特征在文件级也很强", "CLASSIFICATION ANALYSIS", 10)
    for i, (label, acc, f1, color) in enumerate([
        ("1D CNN", 0.9415, 0.9180, COLORS["blue"]),
        ("SVM", 0.9300, 0.9046, COLORS["teal"]),
        ("KNN", 0.9277, 0.9155, COLORS["amber"]),
        ("决策树", 0.8738, 0.8309, COLORS["red"]),
    ]):
        y = Inches(1.55 + i * 0.64)
        bar(s, Inches(0.8), y, label, acc, 1.0, color)
        bar(s, Inches(6.7), y, "Macro-F1", f1, 1.0, color)
    add_image(s, HERE / "figures/experiments/stage2_curves/accuracy_curves.png", Inches(0.8), Inches(4.35), Inches(5.3), Inches(2.1))
    add_image(s, HERE / "figures/experiments/stage2_curves/f1_curves.png", Inches(6.85), Inches(4.35), Inches(5.3), Inches(2.1))
    add_footer(s)

    # 11
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage1 结果：双通道可学习距离信息，但回归精度仍是短板", "REGRESSION RESULT", 11)
    metric_card(s, Inches(0.75), Inches(1.5), "CNN 片段 MAE", "1.90 m", "RMSE 2.43 m", COLORS["blue"])
    metric_card(s, Inches(3.15), Inches(1.5), "文件级 MAE", "2.46 m", "RMSE 2.94 m", COLORS["teal"])
    metric_card(s, Inches(5.55), Inches(1.5), "P95 误差", "4.99 m", "Bias +0.45 m", COLORS["amber"])
    add_image(s, HERE / "figures/experiments/stage1_segment_pred_vs_target.png", Inches(0.8), Inches(2.9), Inches(5.2), Inches(3.55))
    add_image(s, HERE / "figures/experiments/stage1_segment_error_histogram.png", Inches(6.55), Inches(2.9), Inches(5.75), Inches(3.0))
    add_footer(s)

    # 12
    s = prs.slides.add_slide(blank)
    add_title(s, "误差分布：定位困难集中在少样本档位和边界距离", "REGRESSION ERROR ANALYSIS", 12)
    add_image(s, HERE / "figures/experiments/stage1_segment_error_by_distance.png", Inches(0.75), Inches(1.35), Inches(6.25), Inches(3.45))
    add_image(s, HERE / "figures/experiments/error_by_distance_boxplot.png", Inches(7.25), Inches(1.35), Inches(5.1), Inches(3.45))
    small_table(s, Inches(1.0), Inches(5.15), [1.2, 0.7, 0.7, 0.7, 0.7, 0.7, 0.7, 0.7, 0.7], [
        ["距离", "-3", "-1", "0", "+1", "+3", "+5", "+6", "+7"],
        ["MAE", "2.83", "4.85", "0.67", "2.14", "2.44", "1.24", "0.61", "0.95"],
        ["<=1m", "13%", "0%", "74%", "9%", "8%", "37%", "85%", "61%"],
    ], size=8)
    add_lines(s, Inches(8.1), Inches(5.08), Inches(3.95), Inches(1.0), [
        "答辩表达",
        "-1 m 样本数仅 23，误差最大；0 m、+6 m 档位较稳定，说明当前问题更像数据覆盖和档位均衡问题。"
    ], 11, COLORS["ink"], 1.25, True)
    add_footer(s)

    # 13
    s = prs.slides.add_slide(blank)
    add_title(s, "消融与基线：定位真正依赖双通道，SVM 是强回归基线", "ABLATION AND BASELINES", 13)
    small_table(s, Inches(0.75), Inches(1.55), [1.7, 1.2, 1.2, 1.2, 1.2], [
        ["Stage1 方法", "MAE", "RMSE", "<=1m", "文件MAE"],
        ["1D CNN", "1.90", "2.43", "35%", "2.46"],
        ["SVM", "1.51", "1.91", "46%", "1.83"],
        ["KNN", "1.76", "2.29", "44%", "2.02"],
        ["决策树", "2.44", "2.97", "34%", "2.23"],
    ], size=9)
    small_table(s, Inches(7.1), Inches(1.55), [1.8, 1.3, 1.3, 1.3], [
        ["通道设置", "MAE", "RMSE", "<=1m"],
        ["双通道", "1.90", "2.43", "35%"],
        ["单通道", "3.88", "4.75", "17%"],
    ], header_color=COLORS["teal"], size=10)
    add_image(s, ROOT / "消融实验设计示意图.png", Inches(0.9), Inches(4.25), Inches(5.3), Inches(2.15))
    add_lines(s, Inches(7.1), Inches(3.95), Inches(4.75), Inches(1.8), [
        "核心结论",
        "单通道 MAE 从 1.90 m 恶化到 3.88 m，说明距离估计的关键不是单路波形强弱，而是左右通道之间的相对响应差异。"
    ], 13, COLORS["ink"], 1.35, True)
    add_footer(s)

    # 14
    s = prs.slides.add_slide(blank)
    add_title(s, "贡献与边界：框架可行，定位精度仍要靠数据和物理约束继续提升", "CONTRIBUTION AND LIMITATION", 14)
    for i, (title, body, color) in enumerate([
        ("数据流程", "从原始 CSV 到固定窗口样本，并用文件级划分控制评测可信度", COLORS["blue"]),
        ("两阶段框架", "分类与回归解耦，避免正常样本距离标签语义混杂", COLORS["teal"]),
        ("模型比较", "CNN、Conformer、传统基线和通道消融共同解释适用边界", COLORS["amber"]),
    ]):
        y = Inches(1.55 + i * 1.15)
        add_rect(s, Inches(0.85), y, Inches(5.4), Inches(0.9), COLORS["white"], COLORS["line"], True)
        add_text(s, Inches(1.08), y + Inches(0.13), Inches(1.25), Inches(0.3), title, 13, color, True)
        add_text(s, Inches(2.35), y + Inches(0.14), Inches(3.55), Inches(0.48), body, 11, COLORS["ink"])
    add_lines(s, Inches(7.1), Inches(1.55), Inches(4.9), Inches(3.3), [
        "当前边界",
        "1. 数据规模与工况覆盖不足，距离档位不均衡",
        "2. 回归 P95 误差仍接近 5 m，边界档位误差明显",
        "3. 文件级回归简单平均不一定提升结果",
        "4. 跨工况、传感器漂移和在线部署还未验证",
    ], 13, COLORS["ink"], 1.42, True)
    add_rect(s, Inches(0.85), Inches(5.6), Inches(11.1), Inches(0.72), COLORS["navy"], None, True)
    add_text(s, Inches(1.1), Inches(5.78), Inches(10.55), Inches(0.3), "答辩建议：主动承认 SVM 回归更强，把论文贡献定位为“流程、框架和可解释评测”，不要把深度模型包装成全面最优。", 13, COLORS["white"], True, PP_ALIGN.CENTER)
    add_footer(s)

    # 15
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, COLORS["navy"])
    add_rect(s, 0, Inches(3.05), W, Inches(0.09), COLORS["amber"])
    add_text(s, 0, Inches(1.7), W, Inches(0.7), "感谢各位老师批评指正", 38, COLORS["white"], True, PP_ALIGN.CENTER)
    add_text(s, 0, Inches(3.55), W, Inches(0.4), "基于双通道时序信号的两阶段管道泄漏检测与定位方法", 17, RGBColor(225, 232, 242), False, PP_ALIGN.CENTER)
    add_text(s, 0, Inches(4.25), W, Inches(0.32), "答辩人：孟达  |  天津大学建筑工程学院  |  船舶与海洋工程", 13, RGBColor(225, 232, 242), False, PP_ALIGN.CENTER)
    add_text(s, 0, Inches(5.25), W, Inches(0.35), "Q & A", 23, COLORS["amber"], True, PP_ALIGN.CENTER)

    return prs


def make_deck() -> Presentation:
    """Build a balanced 18-slide defense deck with fewer images per slide."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1. Cover
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, COLORS["paper"])
    add_rect(s, 0, 0, W, Inches(0.12), COLORS["amber"])
    add_brand_corner(s)
    add_text(s, Inches(0.92), Inches(0.75), Inches(6.2), Inches(0.3), "天津大学 2026 届本科生毕业设计答辩", 13, COLORS["navy"], True)
    add_text(s, Inches(0.92), Inches(2.15), Inches(10.7), Inches(1.25), "基于双通道时序信号的\n两阶段管道泄漏检测与定位方法", 34, COLORS["navy"], True)
    add_rect(s, Inches(0.95), Inches(3.82), Inches(4.5), Inches(0.045), COLORS["blue"])
    add_text(s, Inches(0.92), Inches(4.28), Inches(9.7), Inches(0.36), "从“是否泄漏”到“泄漏位置”：分类先行，双通道回归定位", 17, COLORS["ink"], True)
    add_lines(s, Inches(0.92), Inches(5.82), Inches(6.9), Inches(0.75), ["答辩人：孟达  |  导师：顾鹏 讲师", "建筑工程学院  |  船舶与海洋工程"], 12, COLORS["muted"], 1.45)
    add_footer(s)

    # 2. Outline
    s = prs.slides.add_slide(blank)
    add_title(s, "汇报结构：围绕“识别可用、定位可行但需继续提升”展开", "OUTLINE", 2)
    outline = [
        ("01", "研究背景", "泄漏检测的工程需求与双通道声信号依据"),
        ("02", "数据与任务", "采集、切分、文件级划分与两阶段标签设计"),
        ("03", "模型与训练", "1D CNN 主干、Conformer 对比与评价体系"),
        ("04", "实验结果", "Stage2 分类、Stage1 回归、基线与消融"),
        ("05", "总结展望", "贡献、局限和下一步改进方向"),
    ]
    for i, (num, title, body) in enumerate(outline):
        y = Inches(1.55 + i * 0.9)
        add_rect(s, Inches(1.0), y, Inches(0.7), Inches(0.42), COLORS["navy"], None, True)
        add_text(s, Inches(1.0), y + Inches(0.08), Inches(0.7), Inches(0.16), num, 10, COLORS["white"], True, PP_ALIGN.CENTER)
        add_text(s, Inches(2.0), y - Inches(0.02), Inches(2.0), Inches(0.3), title, 17, COLORS["navy"], True)
        add_text(s, Inches(4.2), y, Inches(7.6), Inches(0.32), body, 14, COLORS["ink"])
    add_footer(s)

    # 3. Background
    s = prs.slides.add_slide(blank)
    add_title(s, "研究背景：管道泄漏检测需要连续、自动、可定位的监测方法", "BACKGROUND", 3)
    add_image(s, ROOT / "管道泄露危害场景示意图.png", Inches(0.75), Inches(1.45), Inches(5.4), Inches(3.75))
    add_lines(s, Inches(6.65), Inches(1.55), Inches(5.65), Inches(3.2), [
        "工程需求",
        "1. 管道系统广泛用于能源、化工和市政基础设施",
        "2. 泄漏可能造成介质损失、设备损伤和环境风险",
        "3. 人工巡检难以满足连续监测与快速响应要求",
        "4. 自动检测方法需要兼顾识别准确性和定位解释性",
        "5. 单纯报警不足以支撑抢修，距离估计能缩小排查范围",
    ], 14, COLORS["ink"], 1.45, True)
    insight_box(s, Inches(0.95), Inches(5.55), Inches(11.1), Inches(0.92), "本文切入点", "选择双通道传感信号作为对象，不直接做复杂工况建模，而是先验证一条可复现实验链路：状态识别是否稳定、距离信息是否可被学习。", COLORS["navy"])
    add_footer(s)

    # 4. Signal mechanism
    s = prs.slides.add_slide(blank)
    add_title(s, "信号依据：泄漏声传播会形成左右通道响应差异", "SIGNAL MECHANISM", 4)
    add_image(s, ROOT / "管道泄露声信号传播与双通道响应差异.png", Inches(0.75), Inches(1.45), Inches(6.05), Inches(3.85))
    add_lines(s, Inches(7.25), Inches(1.55), Inches(4.9), Inches(3.25), [
        "可学习信息",
        "1. 泄漏点流体喷射引起湍流和管壁振动",
        "2. 扰动沿管壁和流体传播到左右传感器",
        "3. 不同距离带来幅值、相位、到达响应差异",
        "4. 双通道输入比单通道更适合距离估计",
        "5. 分类依赖局部波形差异，定位更依赖左右响应关系",
    ], 14, COLORS["ink"], 1.45, True)
    insight_box(s, Inches(1.0), Inches(5.65), Inches(11.25), Inches(0.82), "核心假设", "双通道原始波形中同时包含两类信息：一类用于判断是否泄漏及泄漏形态，另一类用于通过左右通道差异估计泄漏点相对距离。", COLORS["teal"])
    add_footer(s)

    # 5. Data acquisition
    s = prs.slides.add_slide(blank)
    add_title(s, "数据采集：实验台产生左右通道 CSV 时序信号", "DATA ACQUISITION", 5)
    add_image(s, ROOT / "实验台与数据采集流程图.png", Inches(0.8), Inches(1.45), Inches(6.2), Inches(3.6))
    metric_card(s, Inches(7.45), Inches(1.55), "采样率", "25.6 kHz", "长时序信号", COLORS["blue"])
    metric_card(s, Inches(9.85), Inches(1.55), "窗口长度", "5 s", "128,000 点/通道", COLORS["teal"])
    add_lines(s, Inches(7.45), Inches(3.15), Inches(4.55), Inches(2.2), [
        "处理流程",
        "1. 原始 CSV 保存左右通道长时序信号",
        "2. 以 5 s 固定窗口生成训练片段",
        "3. 对每段信号独立标准化，减少量纲影响",
        "4. 分类样本和回归样本分开组织，避免标签混用",
    ], 14, COLORS["ink"], 1.4, True)
    insight_box(s, Inches(7.45), Inches(5.45), Inches(4.55), Inches(0.78), "数据处理原则", "保留原始波形特征，同时把每个训练样本统一到可比较的长度和尺度。", COLORS["blue"])
    add_footer(s)

    # 6. Dataset design
    s = prs.slides.add_slide(blank)
    add_title(s, "数据划分：文件级划分比片段级随机划分更可信", "DATA SPLIT", 6)
    add_image(s, ROOT / "文件级数据划分示意图.png", Inches(0.75), Inches(1.45), Inches(5.9), Inches(3.65))
    small_table(s, Inches(7.1), Inches(1.65), [1.45, 1.15, 1.15, 1.15, 1.15], [
        ["任务", "总样本", "训练", "验证", "测试"],
        ["Stage2", "11704", "7940", "2464", "1300"],
        ["Stage1", "3737", "2532", "758", "447"],
    ], size=9)
    add_lines(s, Inches(7.1), Inches(3.35), Inches(4.8), Inches(1.6), [
        "关键约束",
        "同一原始文件的全部片段只属于一个数据子集，避免相邻片段同时出现在训练集和测试集。"
    ], 13, COLORS["ink"], 1.35, True)
    add_lines(s, Inches(7.1), Inches(5.0), Inches(4.8), Inches(0.82), [
        "因此，文件级结果比随机片段划分更能反映模型遇到新采集文件时的真实表现。"
    ], 12, COLORS["navy"], 1.25)
    add_footer(s)

    # 7. Two-stage task
    s = prs.slides.add_slide(blank)
    add_title(s, "任务设计：分类和定位目标不同，拆成两阶段更清晰", "TWO-STAGE TASK", 7)
    add_image(s, ROOT / "stage1 与 stage2 标签映射示意.png", Inches(0.85), Inches(1.45), Inches(5.85), Inches(3.5))
    small_table(s, Inches(7.1), Inches(1.7), [1.2, 2.15, 1.65, 1.65], [
        ["阶段", "输入", "输出", "主指标"],
        ["Stage2", "单通道片段", "5 类状态", "Acc / F1"],
        ["Stage1", "左右通道配对", "距离 m", "MAE / RMSE"],
    ], size=9)
    add_text(s, Inches(7.25), Inches(4.0), Inches(4.65), Inches(0.65), "设计理由：正常样本没有距离标签，先分类再回归能避免标签语义混杂。", 14, COLORS["navy"], True)
    insight_box(s, Inches(7.1), Inches(4.9), Inches(4.75), Inches(1.08), "任务边界", "Stage2 负责回答“属于哪种状态”，Stage1 只在泄漏样本上回答“离参考点多远”。这样的拆分更符合物理语义，也便于分别诊断错误来源。", COLORS["teal"])
    add_footer(s)

    # 8. CNN
    s = prs.slides.add_slide(blank)
    add_title(s, "模型一：1D CNN 用多尺度卷积提取长时序波形特征", "1D CNN", 8)
    add_image(s, ROOT / "一维 CNN 网络结构图.png", Inches(0.75), Inches(1.45), Inches(5.45), Inches(3.8))
    add_lines(s, Inches(6.75), Inches(1.55), Inches(5.35), Inches(3.2), [
        "结构要点",
        "1. 四级卷积逐步降采样，将 128000 点压缩为短序列表示",
        "2. 卷积层负责提取局部冲击、持续能量和多尺度模式",
        "3. 全局平均池化降低参数量，减少过拟合风险",
        "4. 分类头与回归头复用同一主干，输出目标不同",
        "5. 对本课题而言，CNN 是可解释、稳定且训练成本较低的强基线",
    ], 13, COLORS["ink"], 1.42, True)
    insight_box(s, Inches(6.75), Inches(5.2), Inches(5.35), Inches(0.82), "为什么不用复杂模型起步", "先用 1D CNN 建立性能基准，可以判断原始波形本身是否足以支撑识别和定位。", COLORS["blue"])
    add_footer(s)

    # 9. Conformer
    s = prs.slides.add_slide(blank)
    add_title(s, "模型二：Conformer 用卷积和自注意力补足长程依赖建模", "CONFORMER", 9)
    add_image(s, ROOT / "conformer 编码块内部结构.png", Inches(0.75), Inches(1.45), Inches(5.9), Inches(3.2))
    add_lines(s, Inches(7.05), Inches(1.55), Inches(4.95), Inches(3.05), [
        "对比目的",
        "1. CNN 更擅长局部波形模式提取",
        "2. Conformer 引入自注意力，补足长程依赖",
        "3. 分类和回归对结构的收益不一定一致",
        "4. 因此本文把 Conformer 作为结构对比，而不是默认最优模型",
    ], 14, COLORS["ink"], 1.4, True)
    small_table(s, Inches(1.05), Inches(5.15), [1.5, 2.1, 2.1, 2.1], [
        ["模型", "优势", "潜在问题", "适用判断"],
        ["CNN", "局部模式稳定", "长程依赖弱", "强基线"],
        ["Conformer", "局部+全局", "参数和超参更敏感", "结构对比"],
    ], size=8)
    add_footer(s)

    # 10. Evaluation
    s = prs.slides.add_slide(blank)
    add_title(s, "评价体系：分类看准确率和 F1，定位看物理误差", "EVALUATION", 10)
    add_image(s, ROOT / "实验评价体系总览.png", Inches(0.75), Inches(1.45), Inches(5.95), Inches(3.7))
    add_lines(s, Inches(7.1), Inches(1.55), Inches(4.95), Inches(3.35), [
        "评价层级",
        "1. 片段级：模型对单个 5 s 窗口的直接判断能力",
        "2. 文件级：同文件多片段聚合，更接近连续采集决策",
        "3. 按距离档位：避免平均值掩盖局部困难",
        "4. 物理误差：MAE、RMSE、Bias、P95、容差命中率",
    ], 13, COLORS["ink"], 1.42, True)
    insight_box(s, Inches(7.1), Inches(5.15), Inches(4.95), Inches(0.86), "评价逻辑", "分类用准确率和 Macro-F1 防止类别不均衡掩盖问题；定位用 MAE、RMSE 和 P95 同时观察平均误差与尾部风险。", COLORS["navy"])
    add_footer(s)

    # 11. Stage2 metrics
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage2 分类：片段级已经较高，文件级聚合更稳定", "CLASSIFICATION METRICS", 11)
    metric_card(s, Inches(0.8), Inches(1.55), "CNN 片段 Acc", "0.9415", "Macro-F1 0.9180", COLORS["blue"])
    metric_card(s, Inches(3.2), Inches(1.55), "CNN 文件 Acc", "0.9783", "Macro-F1 0.9836", COLORS["teal"])
    metric_card(s, Inches(5.6), Inches(1.55), "Conformer 片段 Acc", "0.9900", "Macro-F1 0.9883", COLORS["amber"])
    small_table(s, Inches(0.9), Inches(3.35), [1.5, 1.5, 1.5, 1.5, 1.5], [
        ["方法", "片段Acc", "片段F1", "文件Acc", "文件F1"],
        ["1D CNN", "0.9415", "0.9180", "0.9783", "0.9836"],
        ["SVM", "0.9300", "0.9046", "1.0000", "1.0000"],
        ["KNN", "0.9277", "0.9155", "0.9783", "0.9751"],
        ["决策树", "0.8738", "0.8309", "0.9565", "0.9600"],
    ], size=9)
    add_lines(s, Inches(8.85), Inches(1.62), Inches(3.2), Inches(3.2), [
        "结果解读",
        "1. 分类任务的类别边界较清晰",
        "2. 文件级聚合能显著降低偶然片段误判",
        "3. 传统特征方法在文件级也很强，说明数据中存在明显判别信息",
        "4. Conformer 片段级最高，证明长程上下文对分类有帮助",
    ], 13, COLORS["ink"], 1.35, True)
    add_footer(s)

    # 12. Stage2 confusion
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage2 混淆矩阵：片段级仍有局部混淆，文件级明显收敛", "CONFUSION MATRIX", 12)
    add_image(s, HERE / "figures/experiments/stage2_segment_confusion_matrix.png", Inches(0.75), Inches(1.45), Inches(5.25), Inches(4.25))
    add_image(s, HERE / "figures/experiments/stage2_file_confusion_matrix.png", Inches(6.5), Inches(1.45), Inches(5.25), Inches(4.25))
    add_lines(s, Inches(1.0), Inches(5.95), Inches(10.75), Inches(0.76), [
        "读图结论：片段级误判主要集中在相近泄漏形态，说明单个窗口仍可能受噪声和局部波动影响。",
        "文件级聚合后仅出现 1 个文件误判，表明连续采集条件下分类结果更稳定。"
    ], 12, COLORS["navy"], 1.18)
    add_footer(s)

    # 13. Stage2 training curves
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage2 训练过程：分类指标随训练收敛并保持稳定", "TRAINING CURVES", 13)
    add_image(s, HERE / "figures/experiments/stage2_curves/accuracy_curves.png", Inches(0.75), Inches(1.45), Inches(5.5), Inches(3.45))
    add_image(s, HERE / "figures/experiments/stage2_curves/f1_curves.png", Inches(6.75), Inches(1.45), Inches(5.5), Inches(3.45))
    add_lines(s, Inches(1.0), Inches(5.25), Inches(10.9), Inches(1.05), [
        "曲线说明：模型不是偶然在测试集上得到高分，而是在训练过程中形成了稳定的验证集表现。",
        "Accuracy 与 Macro-F1 同步提升，说明模型并非只偏向多数类；验证曲线后期保持平稳，过拟合风险相对可控。"
    ], 12, COLORS["ink"], 1.22)
    add_footer(s)

    # 14. Stage1 regression
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage1 回归：双通道波形包含距离信息，但精度仍有限", "REGRESSION RESULT", 14)
    metric_card(s, Inches(0.8), Inches(1.55), "片段级 MAE", "1.90 m", "RMSE 2.43 m", COLORS["blue"])
    metric_card(s, Inches(3.2), Inches(1.55), "文件级 MAE", "2.46 m", "RMSE 2.94 m", COLORS["teal"])
    metric_card(s, Inches(5.6), Inches(1.55), "P95 误差", "4.99 m", "Bias +0.45 m", COLORS["amber"])
    add_image(s, HERE / "figures/experiments/stage1_segment_pred_vs_target.png", Inches(0.85), Inches(3.0), Inches(4.8), Inches(3.2))
    add_lines(s, Inches(6.25), Inches(3.1), Inches(5.6), Inches(2.65), [
        "结论",
        "1. 模型已经学到距离变化趋势",
        "2. 预测仍存在向中间距离收缩的现象",
        "3. 文件级平均没有显著优于片段级，说明回归误差中存在系统偏差",
        "4. 定位任务比状态分类更依赖数据覆盖和物理约束",
    ], 14, COLORS["ink"], 1.4, True)
    insight_box(s, Inches(6.25), Inches(5.55), Inches(5.6), Inches(0.72), "答辩表达", "分类解决“能不能发现问题”，回归进一步回答“问题大致在哪里”，两者难度不同。", COLORS["teal"])
    add_footer(s)

    # 15. Stage1 error
    s = prs.slides.add_slide(blank)
    add_title(s, "Stage1 误差：困难集中在少样本档位和边界距离", "ERROR BY DISTANCE", 15)
    add_image(s, HERE / "figures/experiments/stage1_segment_error_by_distance.png", Inches(0.75), Inches(1.45), Inches(6.15), Inches(3.45))
    small_table(s, Inches(7.25), Inches(1.62), [0.9, 0.62, 0.62, 0.62, 0.62, 0.62, 0.62, 0.62, 0.62], [
        ["距离", "-3", "-1", "0", "+1", "+3", "+5", "+6", "+7"],
        ["MAE", "2.83", "4.85", "0.67", "2.14", "2.44", "1.24", "0.61", "0.95"],
        ["<=1m", "13%", "0%", "74%", "9%", "8%", "37%", "85%", "61%"],
    ], size=7)
    add_lines(s, Inches(7.25), Inches(3.35), Inches(4.75), Inches(1.85), [
        "解读",
        "1. -1 m 样本数少且误差最大",
        "2. 0 m 和 +6 m 档位较稳定",
        "3. 边界和少样本档位更容易被预测到中间区域",
        "4. 后续需要补充不均衡档位数据并加入距离约束",
    ], 13, COLORS["ink"], 1.35, True)
    add_footer(s)

    # 16. Baseline and ablation
    s = prs.slides.add_slide(blank)
    add_title(s, "基线与消融：SVM 是强回归基线，双通道是定位关键", "BASELINE / ABLATION", 16)
    small_table(s, Inches(0.75), Inches(1.55), [1.65, 1.2, 1.2, 1.2, 1.25], [
        ["Stage1 方法", "MAE", "RMSE", "<=1m", "文件MAE"],
        ["1D CNN", "1.90", "2.43", "35%", "2.46"],
        ["SVM", "1.51", "1.91", "46%", "1.83"],
        ["KNN", "1.76", "2.29", "44%", "2.02"],
        ["决策树", "2.44", "2.97", "34%", "2.23"],
    ], size=9)
    small_table(s, Inches(7.25), Inches(1.55), [1.7, 1.2, 1.2, 1.2], [
        ["通道设置", "MAE", "RMSE", "<=1m"],
        ["双通道", "1.90", "2.43", "35%"],
        ["单通道", "3.88", "4.75", "17%"],
    ], header_color=COLORS["teal"], size=10)
    add_lines(s, Inches(7.25), Inches(3.35), Inches(4.7), Inches(1.75), [
        "答辩重点",
        "1. 深度模型并非在所有任务上最优",
        "2. SVM 回归更强，说明当前数据规模下传统特征仍有效",
        "3. 单通道误差明显变差，证明定位依赖左右通道差异",
        "4. 本文更重要的贡献是建立可信数据流程、任务拆分和评测框架",
    ], 13, COLORS["ink"], 1.35, True)
    add_footer(s)

    # 17. Contributions and limitations
    s = prs.slides.add_slide(blank)
    add_title(s, "总结：分类任务成熟，定位任务可行但需要继续优化", "CONCLUSION", 17)
    add_lines(s, Inches(0.95), Inches(1.55), Inches(5.55), Inches(4.2), [
        "主要贡献",
        "1. 构建原始 CSV 到固定窗口样本的数据处理流程",
        "2. 采用文件级划分，提升评测可信度",
        "3. 建立 Stage2 分类 + Stage1 回归的两阶段框架",
        "4. 对比 CNN、Conformer、传统基线和通道消融",
        "5. 用片段级、文件级和距离档位共同解释模型表现",
    ], 14, COLORS["ink"], 1.45, True)
    add_lines(s, Inches(7.0), Inches(1.55), Inches(5.0), Inches(4.2), [
        "当前局限",
        "1. 距离档位和工况覆盖仍不足",
        "2. Stage1 回归 P95 误差仍较大",
        "3. 文件级回归聚合策略较简单",
        "4. 跨工况泛化和在线部署尚未验证",
        "5. 物理传播规律尚未显式约束模型输出",
    ], 14, COLORS["ink"], 1.45, True)
    add_footer(s)

    # 18. Main findings
    s = prs.slides.add_slide(blank)
    add_title(s, "主要发现：分类更成熟，定位已验证但仍需增强", "MAIN FINDINGS", 18)
    findings = [
        ("分类任务", "Conformer 片段级 Acc 0.9900，CNN 文件级 Acc 0.9783，说明状态识别已经较稳定。", COLORS["blue"]),
        ("定位任务", "CNN 片段级 MAE 1.90 m，证明双通道波形包含距离信息，但 P95 误差仍偏大。", COLORS["teal"]),
        ("模型边界", "SVM 在回归任务中优于深度模型，说明当前数据规模下传统特征仍有竞争力。", COLORS["amber"]),
        ("关键变量", "单通道 MAE 恶化到 3.88 m，定位依赖左右通道响应差异。", COLORS["red"]),
    ]
    for i, (title, body, color) in enumerate(findings):
        y = Inches(1.45 + i * 1.12)
        add_rect(s, Inches(0.9), y, Inches(1.35), Inches(0.55), color, None, True)
        add_text(s, Inches(0.9), y + Inches(0.14), Inches(1.35), Inches(0.22), title, 11, COLORS["white"], True, PP_ALIGN.CENTER)
        add_text(s, Inches(2.55), y + Inches(0.05), Inches(8.9), Inches(0.45), body, 13, COLORS["ink"])
    add_rect(s, Inches(0.95), Inches(6.15), Inches(10.9), Inches(0.55), COLORS["white"], COLORS["line"], True)
    add_text(s, Inches(1.15), Inches(6.3), Inches(10.5), Inches(0.24), "答辩表达重点：本文不是简单追求某个模型最高分，而是建立一套可信的泄漏识别与定位评测流程。", 12, COLORS["navy"], True, PP_ALIGN.CENTER)
    add_footer(s)

    # 19. Limitation analysis
    s = prs.slides.add_slide(blank)
    add_title(s, "局限分析：Stage1 回归误差来自数据覆盖与物理约束不足", "LIMITATIONS", 19)
    add_image(s, HERE / "figures/experiments/error_histogram.png", Inches(0.75), Inches(1.45), Inches(5.9), Inches(3.25))
    add_lines(s, Inches(7.05), Inches(1.5), Inches(4.95), Inches(3.4), [
        "误差来源判断",
        "1. 部分距离档位样本数量少，模型容易向高频档位收缩",
        "2. 简单文件级平均不一定改善回归，可能放大系统偏差",
        "3. 当前模型主要从数据中学习映射，物理传播约束尚未显式加入",
        "4. 跨压力、流量、管径等工况泛化仍需要验证",
    ], 13, COLORS["ink"], 1.4, True)
    insight_box(s, Inches(7.05), Inches(5.18), Inches(4.95), Inches(0.88), "对结果的态度", "Stage1 的价值在于证明“距离可学习”，不是宣称已经达到工程定位精度。后续优化方向也因此更明确。", COLORS["red"])
    add_footer(s)

    # 20. Future work
    s = prs.slides.add_slide(blank)
    add_title(s, "后续工作：从数据、模型和工程部署三个方向推进", "FUTURE WORK", 20)
    for i, (title, body, color) in enumerate([
        ("数据扩充", "补充少样本距离档位，覆盖更多压力、流量、管径和泄漏强度，并保持文件级独立划分。", COLORS["blue"]),
        ("模型改进", "引入时频融合、距离感知损失和不确定性估计，减少边界档位偏差与中间收缩。", COLORS["teal"]),
        ("决策策略", "改进文件级聚合方式，结合置信度、异常持续时间和报警阈值形成更稳健的报警逻辑。", COLORS["amber"]),
        ("工程验证", "开展跨批次、跨工况测试，评估传感器漂移、噪声变化和在线部署稳定性。", COLORS["red"]),
    ]):
        x = Inches(0.9 + (i % 2) * 5.75)
        y = Inches(1.65 + (i // 2) * 2.05)
        add_rect(s, x, y, Inches(4.95), Inches(1.35), COLORS["white"], COLORS["line"], True)
        add_text(s, x + Inches(0.22), y + Inches(0.2), Inches(4.5), Inches(0.3), title, 16, color, True)
        add_text(s, x + Inches(0.22), y + Inches(0.62), Inches(4.45), Inches(0.5), body, 12, COLORS["ink"])
    add_footer(s)

    # 21. Final takeaway
    s = prs.slides.add_slide(blank)
    add_title(s, "最终结论：两阶段框架可行，下一步重点是提升定位鲁棒性", "TAKEAWAY", 21)
    add_lines(s, Inches(1.1), Inches(1.8), Inches(10.9), Inches(3.0), [
        "本文完成了从原始双通道 CSV 信号到泄漏识别与距离估计的完整流程验证。",
        "Stage2 分类结果表明，泄漏状态识别已经具备较高稳定性；Stage1 回归结果表明，双通道时序信号中确实存在可学习的距离信息。",
        "同时，回归任务的基线对比和误差分析说明，定位精度仍受样本规模、距离档位分布和物理约束不足影响，是后续工作的核心方向。",
        "因此，本文的结论可以概括为：分类链路已经比较成熟，定位链路已经验证可行，工程化还需要进一步提升鲁棒性。",
    ], 18, COLORS["ink"], 1.45)
    add_rect(s, Inches(1.15), Inches(5.45), Inches(10.65), Inches(0.72), COLORS["navy"], None, True)
    add_text(s, Inches(1.35), Inches(5.65), Inches(10.25), Inches(0.26), "一句话概括：分类已经稳定，定位已经可学，工程化还需要更充分的数据和更强的物理约束。", 14, COLORS["white"], True, PP_ALIGN.CENTER)
    add_footer(s)

    # 22. Thanks
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, W, H, COLORS["paper"])
    add_rect(s, 0, 0, W, Inches(0.12), COLORS["amber"])
    add_brand_corner(s)
    add_text(s, Inches(0.95), Inches(1.55), Inches(10.4), Inches(0.8), "感谢各位老师批评指正", 38, COLORS["navy"], True)
    add_rect(s, Inches(0.98), Inches(2.72), Inches(3.65), Inches(0.055), COLORS["blue"])
    add_text(s, Inches(0.98), Inches(3.35), Inches(8.8), Inches(0.35), "基于双通道时序信号的两阶段管道泄漏检测与定位方法", 14, COLORS["ink"])
    add_lines(s, Inches(0.98), Inches(4.18), Inches(6.8), Inches(0.75), ["答辩人：孟达  |  导师：顾鹏 讲师", "天津大学建筑工程学院  |  船舶与海洋工程"], 13, COLORS["muted"], 1.45)
    add_rect(s, Inches(0.98), Inches(5.55), Inches(10.7), Inches(0.72), COLORS["navy"], None, True)
    add_text(s, Inches(1.18), Inches(5.74), Inches(10.25), Inches(0.28), "Q & A  ·  欢迎提问与指导", 16, COLORS["white"], True, PP_ALIGN.CENTER)
    add_footer(s)

    return prs


def main() -> None:
    prs = make_deck()
    prs.save(OUT)
    print(f"written: {OUT}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
