#!/usr/bin/env python3
"""Generate graduation defense PPT for pipeline leak detection thesis."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

PALETTE = {
    "dark_blue": RGBColor(0x1B, 0x2A, 0x4A),
    "blue": RGBColor(0x2B, 0x5C, 0x8A),
    "light_blue": RGBColor(0x3A, 0x7C, 0xBF),
    "accent": RGBColor(0xE8, 0x6A, 0x17),
    "accent2": RGBColor(0x27, 0xAE, 0x60),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "light_gray": RGBColor(0xF2, 0xF4, 0xF7),
    "gray": RGBColor(0xBD, 0xBD, 0xBD),
    "dark_gray": RGBColor(0x33, 0x33, 0x33),
    "subtle": RGBColor(0x66, 0x66, 0x66),
}

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank

def add_bg(slide, color=PALETTE["white"]):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, w, h, text, font_size=18, color=None,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei',
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or PALETTE["dark_gray"]
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox

def add_multiline_textbox(slide, left, top, w, h, lines, font_size=14,
                          color=None, bold_first=False, line_spacing=1.3,
                          alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """Add textbox with multiple paragraphs (lines is list of strings)."""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or PALETTE["dark_gray"]
        p.font.name = font_name
        p.font.bold = (bold_first and i == 0)
        p.alignment = alignment
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_bullet_points(slide, left, top, w, h, items, font_size=14, color=None,
                      bullet_char='●', indent_level=0):
    """Add bulleted text items."""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or PALETTE["dark_gray"]
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        p.line_spacing = Pt(font_size * 1.4)
        p.level = indent_level
    return txBox

def add_section_header(slide, title, subtitle=None):
    """Add a consistent section header bar at top."""
    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PALETTE["accent"])
    # Title area
    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11), Inches(0.7),
                title, font_size=30, color=PALETTE["dark_blue"], bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.7), Inches(0.9), Inches(11), Inches(0.4),
                    subtitle, font_size=14, color=PALETTE["subtle"])
    # Separator line
    add_rect(slide, Inches(0.7), Inches(1.2), Inches(1.5), Inches(0.03), PALETTE["accent"])

def add_page_number(slide, num, total):
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                f"{num}/{total}", font_size=10, color=PALETTE["gray"],
                alignment=PP_ALIGN.RIGHT)

def add_bottom_bar(slide):
    add_rect(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), PALETTE["dark_blue"])

TOTAL_SLIDES = 20

# ==============================================================
# Slide 1: Title
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])

# Large dark blue block on left
add_rect(slide, Inches(0), Inches(0), Inches(6), SLIDE_H, PALETTE["dark_blue"])
# Accent stripe
add_rect(slide, Inches(5.8), Inches(0), Inches(0.2), SLIDE_H, PALETTE["accent"])

# Title on left panel
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(1.8),
            "基于双通道时序信号的\n两阶段管道泄漏检测\n与定位方法",
            font_size=32, color=PALETTE["white"], bold=True)

# Subtitle
add_textbox(slide, Inches(0.8), Inches(3.6), Inches(4.5), Inches(0.5),
            "天津大学 2026 届本科生毕业设计答辩",
            font_size=16, color=PALETTE["gray"])

# Info on right panel
info_lines = [
    "答辩人：孟达",
    "学  院：建筑工程学院",
    "专  业：船舶与海洋工程",
    "学  号：3022205024",
    "导  师：顾鹏 讲师",
]
add_multiline_textbox(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.0),
                      info_lines, font_size=18, color=PALETTE["dark_blue"],
                      line_spacing=1.8)

add_bottom_bar(slide)

# ==============================================================
# Slide 2: Table of Contents
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "汇报大纲")

toc_items = [
    ("01", "研究背景与意义", "管道泄漏检测的工程需求与技术挑战"),
    ("02", "研究现状与问题提出", "传统方法 → 机器学习 → 深度学习 → 本文切入点"),
    ("03", "数据与任务设计", "管道原理、数据集构建、两阶段任务拆分"),
    ("04", "模型设计", "一维 CNN 主干网络、Conformer 结构、训练策略"),
    ("05", "实验结果与分析", "分类与回归结果、基线与消融实验、误差分析"),
    ("06", "总结与展望", "主要贡献、不足与后续改进方向"),
]

y_start = 1.6
for i, (num, title, desc) in enumerate(toc_items):
    y = y_start + i * 0.85
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALETTE["dark_blue"]
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = PALETTE["white"]
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

    # Title + desc
    add_textbox(slide, Inches(1.7), Inches(y - 0.05), Inches(9), Inches(0.35),
                title, font_size=20, color=PALETTE["dark_blue"], bold=True)
    add_textbox(slide, Inches(1.7), Inches(y + 0.3), Inches(9), Inches(0.3),
                desc, font_size=13, color=PALETTE["subtle"])

    # Dotted separator
    if i < len(toc_items) - 1:
        add_rect(slide, Inches(1.7), Inches(y + 0.72), Inches(10), Inches(0.005),
                 PALETTE["light_gray"])

add_bottom_bar(slide)
add_page_number(slide, 2, TOTAL_SLIDES)

# ==============================================================
# Slide 3: Research Background
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "研究背景与意义", "Research Background")

# Left panel - problem description
left_items = [
    "●  管道输送系统广泛应用于能源、化工与市政基础设施",
    "●  泄漏造成介质损失、环境污染、火灾爆炸等次生灾害",
    "●  传统人工巡检难满足连续监测与快速定位需求",
    "●  急需基于传感信号的自动泄漏识别与定位方法",
]
add_multiline_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(2.5),
                      left_items, font_size=16, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Key challenge boxes
challenges = [
    ("复杂噪声", "工况多变，\n背景噪声干扰大"),
    ("长时序信号", "采样率 25600 Hz\n5秒含 128000 点"),
    ("微小泄漏", "非突发泄漏信号\n微弱，难区分"),
    ("精确定位", "传统方法依赖\n传播速度标定"),
]

for i, (title, desc) in enumerate(challenges):
    x = Inches(0.7 + i * 3.1)
    y = Inches(3.8)
    add_rect(slide, x, y, Inches(2.8), Inches(2.5), PALETTE["light_gray"])
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.4),
                title, font_size=16, color=PALETTE["accent"], bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.4), Inches(1.3),
                desc, font_size=13, color=PALETTE["dark_gray"],
                alignment=PP_ALIGN.CENTER)

# Right side key message
add_rect(slide, Inches(10.2), Inches(1.5), Inches(2.6), Inches(2.0), PALETTE["dark_blue"])
add_textbox(slide, Inches(10.4), Inches(1.8), Inches(2.2), Inches(1.5),
            "核心思路\n\n从原始双通道\n波形信号中\n端到端学习\n泄漏特征",
            font_size=14, color=PALETTE["white"], bold=True,
            alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 3, TOTAL_SLIDES)

# ==============================================================
# Slide 4: Literature Review
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "国内外研究现状", "从传统物理模型到深度时序建模的方法演进")

# Three columns
col_data = [
    ("传统物理模型方法", PALETTE["subtle"], [
        "声波法：利用传播时间差",
        "负压波法：压力瞬态检测",
        "流量平衡法：进出口比对",
        "├ 物理解释性强",
        "└ 依赖阈值、抗噪差",
    ]),
    ("机器学习方法", PALETTE["blue"], [
        "SVM / KNN / 随机森林",
        "手工特征 + 分类器组合",
        "时域统计 + 频域能量",
        "├ 可自动学习判别边界",
        "└ 特征设计依赖经验",
    ]),
    ("深度学习方法", PALETTE["accent"], [
        "1D-CNN：局部波形特征",
        "Transformer：长程依赖",
        "Conformer：卷积+注意力",
        "├ 端到端减少人工设计",
        "└ 需充分数据与防泄漏划分",
    ]),
]

for i, (title, accent_color, items) in enumerate(col_data):
    x = Inches(0.5 + i * 4.2)
    # Header
    add_rect(slide, x, Inches(1.5), Inches(3.9), Inches(0.55), accent_color)
    add_textbox(slide, x + Inches(0.1), Inches(1.55), Inches(3.7), Inches(0.45),
                title, font_size=15, color=PALETTE["white"], bold=True,
                alignment=PP_ALIGN.CENTER)
    # Items
    add_multiline_textbox(slide, x + Inches(0.15), Inches(2.2), Inches(3.6), Inches(3.5),
                          items, font_size=13, color=PALETTE["dark_gray"],
                          line_spacing=1.6)

# Bottom arrow
add_textbox(slide, Inches(0.7), Inches(5.4), Inches(11.8), Inches(0.5),
            "▸ 发展趋势：从物理模型驱动 → 数据驱动，从人工特征 → 端到端表示学习，从单一任务 → 多任务协同",
            font_size=13, color=PALETTE["accent"], bold=True)

# Key gaps
gaps = [
    "① 传统方法：特征设计依赖人工，泛化能力有限，定位精度不足",
    "② 现有深度方法：大多单任务处理，分类与回归目标差异未被充分利用",
    "③ 本文切入点：两阶段先分类后定位 + 1D-CNN/Conformer 结构对比 + 多层级评测",
]
add_multiline_textbox(slide, Inches(0.7), Inches(5.7), Inches(11.8), Inches(1.5),
                      gaps, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

add_bottom_bar(slide)
add_page_number(slide, 4, TOTAL_SLIDES)

# ==============================================================
# Slide 5: Research Problem & Work
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "研究问题与本文工作", "Research Problem & Contributions")

# Core problem
add_rect(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.6), PALETTE["light_gray"])
add_textbox(slide, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.5),
            "核心问题：如何利用双通道长时序传感信号，在减少人工特征设计的前提下，实现管道泄漏状态识别与泄漏距离估计？",
            font_size=16, color=PALETTE["dark_blue"], bold=True)

# Three aspects
aspects = [
    ("数据层面", "原始 CSV → 固定窗口切分\n→ 文件级划分避免信息泄漏"),
    ("任务层面", "Stage2 分类：状态识别\nStage1 回归：距离估计"),
    ("模型层面", "1D CNN 基础模型\n+ Conformer 结构对比"),
]
for i, (title, desc) in enumerate(aspects):
    x = Inches(0.7 + i * 4.1)
    add_rect(slide, x, Inches(2.4), Inches(3.8), Inches(1.6), PALETTE["light_gray"])
    add_textbox(slide, x + Inches(0.15), Inches(2.5), Inches(3.5), Inches(0.35),
                title, font_size=16, color=PALETTE["accent"], bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(2.9), Inches(3.5), Inches(1.0),
                desc, font_size=14, color=PALETTE["dark_gray"])

# Main work items
work_items = [
    "① 构建了原始 CSV 信号 → 固定窗口训练样本的自动化处理流程，采用文件级划分避免数据泄漏",
    "② 建立两阶段学习框架：Stage2 分类（单通道 5分类）+ Stage1 回归（双通道距离估计）",
    "③ 实现 1D CNN 基础模型（4层带步长卷积）+ Conformer 结构（卷积+自注意力混合）",
    "④ 搭建统一训练框架，补充传统 ML 基线、通道消融和结构对比实验，多层级评测",
]
add_multiline_textbox(slide, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.5),
                      work_items, font_size=14, color=PALETTE["dark_gray"],
                      line_spacing=1.8)

add_bottom_bar(slide)
add_page_number(slide, 5, TOTAL_SLIDES)

# ==============================================================
# Slide 6: Pipeline Principle
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "管道泄漏探测原理与信号特征", "Pipeline Leak Detection Principle")

# Left: physical principle
principles = [
    "● 泄漏点流体喷射 → 湍流、压力脉动、管壁振动",
    "● 扰动沿管壁和流体传播 → 左右传感器接收响应差异",
    "● 泄漏点距传感器越近 → 幅值衰减越小、响应越早",
    "● 左右通道幅值差异、相位偏移 → 距离估计物理依据",
]
add_multiline_textbox(slide, Inches(0.7), Inches(1.5), Inches(6.0), Inches(2.2),
                      principles, font_size=14, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Signal characteristics
add_textbox(slide, Inches(0.7), Inches(3.5), Inches(6.0), Inches(0.4),
            "信号关键特征", font_size=16, color=PALETTE["dark_blue"], bold=True)

sig_features = [
    "• 采样率 25,600 Hz × 5 s = 128,000 点/通道",
    "• 非平稳性：泄漏扰动随工况变化，片段间差异大",
    "• 双通道互补：单侧信号难区分「距离近」与「强度大」",
    "• 多尺度模式：局部冲击 + 持续能量变化 + 周期性振荡",
]
add_multiline_textbox(slide, Inches(0.7), Inches(3.9), Inches(6.0), Inches(2.2),
                      sig_features, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.6)

# Right: data flow
add_rect(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.0), PALETTE["light_gray"])
flow_steps = [
    "① 实验台采集 → 左右通道 CSV 文件",
    "② 固定窗口切分（5 s 非重叠）",
    "③ 逐段标准化 (z-score)",
    "④ 文件级划分 → 训练/验证/测试集",
    "⑤ 分类：左右独立单通道样本",
    "⑥ 回归：同窗口双通道配对样本",
]
add_multiline_textbox(slide, Inches(7.5), Inches(1.7), Inches(5.0), Inches(4.5),
                      flow_steps, font_size=14, color=PALETTE["dark_gray"],
                      line_spacing=1.9)

add_bottom_bar(slide)
add_page_number(slide, 6, TOTAL_SLIDES)

# ==============================================================
# Slide 7: Dataset Construction
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "数据集构建与任务拆分", "Dataset & Task Design")

# Dataset stats
add_rect(slide, Inches(0.7), Inches(1.5), Inches(3.5), Inches(1.0), PALETTE["dark_blue"])
add_textbox(slide, Inches(0.9), Inches(1.6), Inches(3.1), Inches(0.7),
            "数据集统计\nStage2: 11,704 样本 | Stage1: 3,737 配对样本",
            font_size=14, color=PALETTE["white"], bold=True,
            alignment=PP_ALIGN.CENTER)

# Task splitting diagram - Two boxes side by side
# Stage2
add_rect(slide, Inches(0.7), Inches(2.9), Inches(5.8), Inches(3.2), PALETTE["light_gray"])
add_textbox(slide, Inches(0.9), Inches(3.0), Inches(2.0), Inches(0.4),
            "Stage2 分类任务", font_size=18, color=PALETTE["accent"], bold=True)

s2_items = [
    "输入：单通道分段信号 (1 × 128,000)",
    "输出：五类标签（正常 + 多类泄漏形态）",
    "损失函数：交叉熵损失 (Cross-Entropy)",
    "评价指标：Accuracy、Macro-F1、混淆矩阵",
    "评测层级：片段级 + 文件级聚合",
    "训练/验证/测试：7,940 / 2,464 / 1,300",
]
add_multiline_textbox(slide, Inches(0.9), Inches(3.5), Inches(5.4), Inches(2.5),
                      s2_items, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.6)

# Stage1
add_rect(slide, Inches(6.9), Inches(2.9), Inches(5.8), Inches(3.2), PALETTE["light_gray"])
add_textbox(slide, Inches(7.1), Inches(3.0), Inches(2.0), Inches(0.4),
            "Stage1 回归任务", font_size=18, color=PALETTE["accent2"], bold=True)

s1_items = [
    "输入：左右双通道配对 (2 × 128,000)",
    "输出：泄漏距离（连续数值，单位 m）",
    "损失函数：Smooth L1 损失",
    "评价指标：MAE、RMSE、容差命中率、Bias",
    "评测层级：片段级 + 文件级 + 按距离档位",
    "训练/验证/测试：2,532 / 758 / 447",
]
add_multiline_textbox(slide, Inches(7.1), Inches(3.5), Inches(5.4), Inches(2.5),
                      s1_items, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.6)

# Key design note
add_rect(slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5), PALETTE["light_gray"])
add_textbox(slide, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.45),
            "⚡ 关键设计：文件级划分 — 同文件的全部片段只属于一个数据子集，避免相邻片段泄漏到训练/测试集导致指标虚高",
            font_size=14, color=PALETTE["accent"], bold=True)

add_bottom_bar(slide)
add_page_number(slide, 7, TOTAL_SLIDES)

# ==============================================================
# Slide 8: Two-Stage Framework
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "两阶段检测框架总体流程", "Two-Stage Detection Pipeline")

# Flow diagram using boxes and arrows
boxes = [
    (0.3, "原始双通道\nCSV 信号", PALETTE["dark_blue"]),
    (3.0, "固定窗口切分\n逐段标准化", PALETTE["blue"]),
    (5.7, "Stage2 分类\n单通道输入", PALETTE["accent"]),
    (8.7, "是否泄漏?", PALETTE["accent"]),
    (8.7, "Stage1 回归\n双通道输入", PALETTE["accent2"]),
    (11.0, "距离估计\n结果输出", PALETTE["accent2"]),
]

for x, label, color in boxes:
    add_rect(slide, Inches(x), Inches(2.0), Inches(2.2), Inches(1.0), color)
    add_textbox(slide, Inches(x + 0.1), Inches(2.1), Inches(2.0), Inches(0.8),
                label, font_size=12, color=PALETTE["white"], bold=True,
                alignment=PP_ALIGN.CENTER)

# Arrows (simple text arrows)
for a_x in [2.6, 5.1, 7.8]:
    add_textbox(slide, Inches(a_x), Inches(2.15), Inches(0.4), Inches(0.6),
                "→", font_size=24, color=PALETTE["dark_blue"], bold=True,
                alignment=PP_ALIGN.CENTER)

# Branch arrows
add_textbox(slide, Inches(8.2), Inches(3.2), Inches(1.5), Inches(0.4),
            "否 → 输出类别", font_size=12, color=PALETTE["subtle"])
add_textbox(slide, Inches(8.2), Inches(2.65), Inches(0.4), Inches(0.5),
            "↓ 是", font_size=12, color=PALETTE["accent2"], bold=True)

# Design rationale
rationale = [
    "设计理由：",
    "① 状态识别与距离估计共享同一信号来源，但学习目标不同（离散判别 vs 连续回归）",
    "② 正常样本只有分类意义，无距离标签 → 两阶段拆分避免标签语义混杂",
    "③ 工程优势：Stage2 常驻监测 → 仅泄漏时触发 Stage1 → 减少无效计算",
    "④ 模块化：两阶段可分别优化、独立更新、设置不同报警阈值",
]
add_multiline_textbox(slide, Inches(0.7), Inches(3.8), Inches(11.9), Inches(2.8),
                      rationale, font_size=14, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

add_bottom_bar(slide)
add_page_number(slide, 8, TOTAL_SLIDES)

# ==============================================================
# Slide 9: 1D CNN Architecture
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "一维卷积特征提取网络 (1D-CNN Backbone)", "1D CNN Architecture")

# Architecture table
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.4),
            "网络结构：4 级卷积逐步降采样 → 全局平均池化 → 全连接任务头",
            font_size=16, color=PALETTE["dark_blue"], bold=True)

# Layer visualization
layers = [
    ("Conv1", "C=32, K=15, S=4", "128000 → 32000", "局部波形\n细节提取"),
    ("Conv2", "C=64, K=9, S=4", "32000 → 8000", "中层特征\n模式组合"),
    ("Conv3", "C=128, K=7, S=4", "8000 → 2000", "时域压缩\n感受野扩大"),
    ("Conv4", "C=256, K=5, S=4", "2000 → 500", "高级语义\n全局表示"),
    ("Pool+FC", "AdaptiveAvgPool", "500 → 256 → 输出", "→ 分类/回归"),
]

for i, (name, params, dim, role) in enumerate(layers):
    y = Inches(2.2 + i * 0.8)
    colors = [PALETTE["dark_blue"], PALETTE["blue"], PALETTE["light_blue"],
              PALETTE["accent"], PALETTE["accent2"]]
    add_rect(slide, Inches(0.7), y, Inches(1.8), Inches(0.6), colors[i])
    add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(1.6), Inches(0.5),
                name, font_size=16, color=PALETTE["white"], bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2.7), y + Inches(0.05), Inches(3.0), Inches(0.5),
                params, font_size=14, color=PALETTE["dark_gray"])
    add_textbox(slide, Inches(5.7), y + Inches(0.05), Inches(3.0), Inches(0.5),
                dim, font_size=14, color=PALETTE["dark_gray"])
    add_textbox(slide, Inches(8.7), y + Inches(0.05), Inches(2.0), Inches(0.5),
                role, font_size=13, color=PALETTE["subtle"])

    if i < len(layers) - 1:
        add_textbox(slide, Inches(1.4), y + Inches(0.5), Inches(0.4), Inches(0.3),
                    "▼", font_size=12, color=PALETTE["gray"], alignment=PP_ALIGN.CENTER)

# Key design notes
notes = [
    "● 每层 = Conv1d + BatchNorm + GELU + Conv1d + BatchNorm + GELU",
    "● 步长 4 逐步降采样：时序长度每次缩至 1/4，感受野逐层扩大",
    "● GELU 激活函数：相比 ReLU，负值区域保留信息，非线性更平滑",
    "● 全局平均池化：减少参数量，降低过拟合风险",
]
add_multiline_textbox(slide, Inches(0.7), Inches(6.2), Inches(11.9), Inches(1.0),
                      notes, font_size=12, color=PALETTE["dark_gray"],
                      line_spacing=1.5)

add_bottom_bar(slide)
add_page_number(slide, 9, TOTAL_SLIDES)

# ==============================================================
# Slide 10: Conformer Architecture
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "Conformer 时序建模结构", "Conformer: CNN + Self-Attention Hybrid")

# Left: Conformer block
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.4),
            "Conformer 编码块结构", font_size=16, color=PALETTE["dark_blue"], bold=True)

conformer_components = [
    "① 卷积前端降采样（同 CNN，128000 → 500）",
    "② 映射到隐空间 (d_model) + 位置编码",
    "③ 堆叠 N 个 Conformer 编码层：",
    "    ├ 前馈网络 (FFN) — 逐位置非线性变换",
    "    ├ 多头自注意力 (MHSA) — 全局依赖建模",
    "    ├ 卷积模块 (Conv) — 局部时序模式",
    "    └ 残差连接 + LayerNorm — 稳定训练",
    "④ 全局池化 → 任务头（同 CNN）",
]
add_multiline_textbox(slide, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.5),
                      conformer_components, font_size=14, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Right: comparison
add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.4),
            "CNN vs Conformer 对比", font_size=16, color=PALETTE["dark_blue"], bold=True)

# CNN column
add_rect(slide, Inches(7.2), Inches(2.1), Inches(2.6), Inches(0.45), PALETTE["blue"])
add_textbox(slide, Inches(7.3), Inches(2.12), Inches(2.4), Inches(0.4),
            "1D CNN", font_size=15, color=PALETTE["white"], bold=True,
            alignment=PP_ALIGN.CENTER)

cnn_pros = [
    "✓ 局部卷积建模",
    "✓ 训练稳定、参数量可控",
    "✓ 梯度传播顺畅",
    "✗ 感受野受层数限制",
    "✗ 长程依赖建模弱",
]
add_multiline_textbox(slide, Inches(7.3), Inches(2.7), Inches(2.4), Inches(2.5),
                      cnn_pros, font_size=12, color=PALETTE["dark_gray"],
                      line_spacing=1.6)

# Conformer column
add_rect(slide, Inches(10.0), Inches(2.1), Inches(2.6), Inches(0.45), PALETTE["accent"])
add_textbox(slide, Inches(10.1), Inches(2.12), Inches(2.4), Inches(0.4),
            "Conformer", font_size=15, color=PALETTE["white"], bold=True,
            alignment=PP_ALIGN.CENTER)

conf_pros = [
    "✓ 局部 + 全局建模",
    "✓ 多头注意力多视角",
    "✓ 分类任务提升显著",
    "✗ 计算复杂度更高",
    "✗ 回归未见一致收益",
]
add_multiline_textbox(slide, Inches(10.1), Inches(2.7), Inches(2.4), Inches(2.5),
                      conf_pros, font_size=12, color=PALETTE["dark_gray"],
                      line_spacing=1.6)

# Only use when CNN前端降采样
add_rect(slide, Inches(7.2), Inches(5.3), Inches(5.5), Inches(1.2), PALETTE["light_gray"])
add_textbox(slide, Inches(7.4), Inches(5.4), Inches(5.1), Inches(1.0),
            "⚡ 关键设计：注意力仅在降采样后的特征序列上计算\n"
            "原始 128000 点 → 卷积前端压缩 → ~500 长度 → 注意力\n"
            "避免 O(L²) 复杂度，兼顾计算可行性与长程建模",
            font_size=12, color=PALETTE["accent"])

add_bottom_bar(slide)
add_page_number(slide, 10, TOTAL_SLIDES)

# ==============================================================
# Slide 11: Training Strategy
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "训练策略与正则化", "Training Strategy & Regularization")

# Optimization
strategies = [
    ("优化器", "AdamW\n(lr_cnn=0.001, lr_conf=0.0005)\n解耦权重衰减 wd=1e-4", PALETTE["dark_blue"]),
    ("学习率调度", "余弦退火 (Cosine Annealing)\nη_t = η_min + ½(η_max-η_min)(1+cos(tπ/T))\n平滑衰减，减少震荡", PALETTE["blue"]),
    ("早停机制", "Patience = 8 epochs\n监控验证集指标\n分类→Acc, 回归→MAE\n恢复到最优模型", PALETTE["light_blue"]),
    ("正则化", "权重衰减 L2 (λ=0.0001)\n梯度裁剪 max_norm=1.0\n批归一化隐式正则\n逐段标准化消除幅值差异", PALETTE["accent"]),
]

for i, (title, desc, color) in enumerate(strategies):
    x = Inches(0.5 + i * 3.15)
    add_rect(slide, x, Inches(1.5), Inches(2.95), Inches(3.0), color)
    add_textbox(slide, x + Inches(0.15), Inches(1.6), Inches(2.65), Inches(0.35),
                title, font_size=16, color=PALETTE["white"], bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(2.1), Inches(2.65), Inches(2.2),
                desc, font_size=12, color=PALETTE["white"],
                alignment=PP_ALIGN.CENTER)

# Training flow
flow_str = "配置加载 → 数据加载 → 模型创建 → Epoch 循环 → [前向 → 损失 → 反向 → 梯度裁剪 → 参数更新] → 验证评估 → 早停判断 → 测试集评测"
add_rect(slide, Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.6), PALETTE["light_gray"])
add_textbox(slide, Inches(0.9), Inches(4.98), Inches(11.5), Inches(0.45),
            f"训练流程：{flow_str}",
            font_size=12, color=PALETTE["dark_gray"])

# Key points
add_textbox(slide, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.4),
            "重要保障措施", font_size=16, color=PALETTE["dark_blue"], bold=True)

safeguards = [
    "● 随机种子固定 (seed=42)：保证数据划分和初始化可复现",
    "● 文件级划分：同一原始文件的片段全部归入同一数据子集",
    "● 分类以验证集 Accuracy 选模型，回归以验证集 MAE 选模型",
    "● 统一框架：配置文件区分任务类型，训练器代码完全复用",
]
add_multiline_textbox(slide, Inches(0.7), Inches(6.2), Inches(11.9), Inches(1.0),
                      safeguards, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.6)

add_bottom_bar(slide)
add_page_number(slide, 11, TOTAL_SLIDES)

# ==============================================================
# Slide 12: Evaluation Metrics
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "实验评价体系", "Evaluation Framework")

# Classification metrics
add_rect(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(2.0), PALETTE["light_gray"])
add_textbox(slide, Inches(0.9), Inches(1.55), Inches(5.4), Inches(0.35),
            "Stage2 分类指标", font_size=16, color=PALETTE["accent"], bold=True)
cls_metrics = [
    "• Accuracy = N_correct / N_total",
    "• Macro-F1：各类别 F1 等权平均，反映少数类识别",
    "• 混淆矩阵：展示类别间误分类模式",
    "• 评测层级：片段级（单窗口判别）+ 文件级（聚合判断）",
]
add_multiline_textbox(slide, Inches(0.9), Inches(2.0), Inches(5.4), Inches(1.4),
                      cls_metrics, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Regression metrics
add_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.0), PALETTE["light_gray"])
add_textbox(slide, Inches(7.1), Inches(1.55), Inches(5.4), Inches(0.35),
            "Stage1 回归指标", font_size=16, color=PALETTE["accent2"], bold=True)
reg_metrics = [
    "• MAE = Σ|d_i - d̂_i| / N  （平均绝对误差）",
    "• RMSE = √(Σ(d_i - d̂_i)²/N)  （对大误差敏感）",
    "• Bias = mean(d̂_i - d_i)  （系统偏差方向）",
    "• 容差命中率：|d̂_i - d_i| ≤ τ 的样本占比",
]
add_multiline_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.4), Inches(1.4),
                      reg_metrics, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Multi-level evaluation
add_textbox(slide, Inches(0.7), Inches(3.8), Inches(11.9), Inches(0.4),
            "多层级评测设计", font_size=16, color=PALETTE["dark_blue"], bold=True)

eval_levels = [
    ("片段级评测", "模型对单个 5 s 窗口的直接判断能力\n反映局部观测条件下的即时表现"),
    ("文件级评测", "同文件多片段聚合后判断\n更接近工程中连续采集的综合决策"),
    ("按距离档位拆分", "不同距离位置的误差分别报告\n避免总体均值掩盖局部困难"),
    ("物理误差分析", "Bias、P95、容差命中率、R²\n判断是否存在系统性偏差或大误差"),
]
for i, (title, desc) in enumerate(eval_levels):
    x = Inches(0.5 + i * 3.15)
    add_rect(slide, x, Inches(4.3), Inches(2.95), Inches(1.5), PALETTE["light_gray"])
    add_textbox(slide, x + Inches(0.1), Inches(4.35), Inches(2.75), Inches(0.3),
                title, font_size=14, color=PALETTE["accent"], bold=True)
    add_textbox(slide, x + Inches(0.1), Inches(4.7), Inches(2.75), Inches(0.9),
                desc, font_size=11, color=PALETTE["dark_gray"])

# Experiment scope
add_rect(slide, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.8), PALETTE["dark_blue"])
exp_scope = [
    "实验矩阵：主模型 × 2（CNN / Conformer）| 任务 × 2（分类 / 回归）| 基线 × 3（SVM / KNN / 决策树）| 消融 × 2（通道 / 损失函数）",
    "传统基线：时域统计特征 + 频域能量特征 → SVM / KNN / 决策树",
]
add_multiline_textbox(slide, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.7),
                      exp_scope, font_size=13, color=PALETTE["white"],
                      line_spacing=1.6)

add_bottom_bar(slide)
add_page_number(slide, 12, TOTAL_SLIDES)

# ==============================================================
# Slide 13: Stage2 Classification Results
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "Stage2 分类实验结果", "Classification Results (5-class)")

# Main results table
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4),
            "一维 CNN 分类结果", font_size=16, color=PALETTE["dark_blue"], bold=True)

# Simple visual table
table_data = [
    ["评测层级", "Accuracy", "Macro-F1"],
    ["片段级", "0.9415", "0.9180"],
    ["文件级", "0.9783", "0.9836"],
]
for row_i, row in enumerate(table_data):
    for col_j, cell in enumerate(row):
        x = Inches(0.7 + col_j * 2.0)
        y = Inches(2.0 + row_i * 0.5)
        if row_i == 0:
            add_rect(slide, x, y, Inches(1.9), Inches(0.45), PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        else:
            add_rect(slide, x, y, Inches(1.9), Inches(0.45), PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.7), Inches(0.35),
                    cell, font_size=14, color=text_color, bold=(row_i == 0),
                    alignment=PP_ALIGN.CENTER)

# Confusion matrix note
add_textbox(slide, Inches(0.7), Inches(3.8), Inches(5.5), Inches(0.4),
            "混淆矩阵分析", font_size=16, color=PALETTE["dark_blue"], bold=True)

conf_analysis = [
    "● 类别 4：识别最稳定，测试集全部正确",
    "● 类别 1 ↔ 3：31 个片段混淆（局部波形相似）",
    "● 类别 2 → 0/1：15 + 27 个片段误判",
    "● 文件级聚合后：46 个文件中仅 1 个类别 2 误判",
]
add_multiline_textbox(slide, Inches(0.7), Inches(4.2), Inches(5.5), Inches(2.0),
                      conf_analysis, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Baselines comparison
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
            "传统基线对比", font_size=16, color=PALETTE["dark_blue"], bold=True)

base_table = [
    ["方法", "片段Acc", "片段F1", "文件Acc", "文件F1"],
    ["1D CNN", "0.9415", "0.9180", "0.9783", "0.9836"],
    ["SVM", "0.9300", "0.9046", "1.0000", "1.0000"],
    ["KNN", "0.9277", "0.9155", "0.9783", "0.9751"],
    ["决策树", "0.8738", "0.8309", "0.9565", "0.9600"],
]

for row_i, row in enumerate(base_table):
    for col_j, cell in enumerate(row):
        x = Inches(7.0 + col_j * 1.35)
        y = Inches(2.0 + row_i * 0.48)
        if row_i == 0:
            add_rect(slide, x, y, Inches(1.3), Inches(0.43), PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        elif row_i == 1:
            add_rect(slide, x, y, Inches(1.3), Inches(0.43), PALETTE["light_gray"])
            text_color = PALETTE["accent"]
        elif row_i == 2:
            add_rect(slide, x, y, Inches(1.3), Inches(0.43), PALETTE["white"])
            text_color = PALETTE["dark_gray"]
        else:
            add_rect(slide, x, y, Inches(1.3), Inches(0.43), PALETTE["white"])
            text_color = PALETTE["dark_gray"]
        add_textbox(slide, x + Inches(0.05), y + Inches(0.05), Inches(1.2), Inches(0.33),
                    cell, font_size=11, color=text_color, bold=(row_i == 0),
                    alignment=PP_ALIGN.CENTER)

# Key insight
add_rect(slide, Inches(7.0), Inches(4.6), Inches(5.7), Inches(1.6), PALETTE["light_gray"])
insight_lines = [
    "💡 关键发现：",
    "① CNN 片段级优于传统方法，端到端学习",
    "   对单片段判别有收益",
    "② 文件级 SVM 达 1.000，说明文件级可分性强",
    "③ 分类数据中统计特征已含明显判别信息",
]
add_multiline_textbox(slide, Inches(7.2), Inches(4.7), Inches(5.3), Inches(1.4),
                      insight_lines, font_size=12, color=PALETTE["dark_gray"],
                      line_spacing=1.6)

add_bottom_bar(slide)
add_page_number(slide, 13, TOTAL_SLIDES)

# ==============================================================
# Slide 14: Conformer Classification Results
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "Stage2 CNN vs Conformer 分类对比", "Architecture Comparison for Classification")

# Big comparison
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.4),
            "一维 CNN vs Conformer — Stage2 五分类结果对比",
            font_size=18, color=PALETTE["dark_blue"], bold=True,
            alignment=PP_ALIGN.CENTER)

# Metrics comparison table
comp_data = [
    ["模型", "片段级 Acc", "片段级 Macro-F1", "文件级 Acc", "文件级 Macro-F1"],
    ["1D CNN", "0.9415", "0.9180", "0.9783", "0.9836"],
    ["Conformer", "0.9900 ▲", "0.9883 ▲", "1.0000 ▲", "1.0000 ▲"],
]

for row_i, row in enumerate(comp_data):
    for col_j, cell in enumerate(row):
        x = Inches(1.0 + col_j * 2.9)
        y = Inches(2.2 + row_i * 0.65)
        w = Inches(2.7)
        h = Inches(0.55)
        if row_i == 0:
            add_rect(slide, x, y, w, h, PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        elif row_i == 1:
            add_rect(slide, x, y, w, h, PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        else:
            add_rect(slide, x, y, w, h, PALETTE["accent"])
            text_color = PALETTE["white"]
        add_textbox(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), h - Inches(0.16),
                    cell, font_size=15, color=text_color,
                    bold=(row_i >= 1), alignment=PP_ALIGN.CENTER)

# Improvement arrows
improvements = [
    "Conformer 片段级 Accuracy 提升: 0.9415 → 0.9900 (+5.2%)",
    "Conformer 片段级 Macro-F1 提升: 0.9180 → 0.9883 (+7.7%)",
    "文件级 Accuracy 和 Macro-F1 均达到 100%",
]
add_multiline_textbox(slide, Inches(1.5), Inches(3.7), Inches(10.0), Inches(1.2),
                      improvements, font_size=15, color=PALETTE["accent"],
                      line_spacing=1.8)

# Analysis
add_textbox(slide, Inches(0.7), Inches(4.8), Inches(11.9), Inches(0.4),
            "分析", font_size=16, color=PALETTE["dark_blue"], bold=True)

analysis = [
    "● 自注意力机制有效捕获长程上下文 → 多个片段互相印证，减少局部噪声干扰",
    "● 卷积模块保留局部波形建模 → 短时冲击和局部振荡不被遗漏",
    "● 注意力在 CNN 降采样后计算 → 计算可行（~500 长度而非 128000）",
    "● Conformer 训练对超参数更敏感（学习率 0.0005 vs CNN 0.001）",
]
add_multiline_textbox(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.8),
                      analysis, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

add_bottom_bar(slide)
add_page_number(slide, 14, TOTAL_SLIDES)

# ==============================================================
# Slide 15: Stage1 Regression Results
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "Stage1 距离回归实验结果", "Distance Regression Results")

# Main results
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4),
            "一维 CNN 回归结果", font_size=16, color=PALETTE["dark_blue"], bold=True)

reg_table = [
    ["评测层级", "MAE (m)", "RMSE (m)"],
    ["片段级", "1.90", "2.43"],
    ["文件级", "2.46", "2.94"],
]
for row_i, row in enumerate(reg_table):
    for col_j, cell in enumerate(row):
        x = Inches(0.7 + col_j * 1.8)
        y = Inches(2.0 + row_i * 0.5)
        if row_i == 0:
            add_rect(slide, x, y, Inches(1.7), Inches(0.45), PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        else:
            add_rect(slide, x, y, Inches(1.7), Inches(0.45), PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.5), Inches(0.35),
                    cell, font_size=14, color=text_color, bold=(row_i == 0),
                    alignment=PP_ALIGN.CENTER)

# Additional metrics
add_textbox(slide, Inches(0.7), Inches(3.1), Inches(5.5), Inches(0.4),
            "物理误差分析", font_size=14, color=PALETTE["dark_blue"], bold=True)

phy_metrics = [
    "中位绝对误差: 1.45 m  |  P95 误差: 4.99 m",
    "最大绝对误差: 6.86 m  |  R² = 0.54",
    "Bias (有符号偏差): +0.45 m (整体偏正)",
    "误差 ≤ 1.0 m: 35%  | 误差 ≤ 2.0 m: 66%",
]
add_multiline_textbox(slide, Inches(0.7), Inches(3.4), Inches(5.5), Inches(2.0),
                      phy_metrics, font_size=12, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Distance-wise analysis
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
            "按距离档位误差分布", font_size=16, color=PALETTE["dark_blue"], bold=True)

dist_table = [
    ["距离 (m)", "-3", "-1", "0", "+1", "+3", "+5", "+6", "+7"],
    ["样本数", "118", "23", "69", "46", "53", "46", "46", "46"],
    ["MAE", "2.83", "4.85", "0.67", "2.14", "2.44", "1.24", "0.61", "0.95"],
    ["≤1m%", "13%", "0%", "74%", "9%", "8%", "37%", "85%", "61%"],
]

for row_i, row in enumerate(dist_table):
    for col_j, cell in enumerate(row):
        x = Inches(7.0 + col_j * 0.7)
        y = Inches(2.0 + row_i * 0.45)
        w = Inches(0.68)
        h = Inches(0.4)
        if row_i == 0:
            add_rect(slide, x, y, w, h, PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        elif "74%" in cell or "85%" in cell:
            add_rect(slide, x, y, w, h, PALETTE["accent2"])
            text_color = PALETTE["white"]
        elif "4.85" in cell or "0%" in cell:
            add_rect(slide, x, y, w, h, RGBColor(0xE7, 0x4C, 0x3C))
            text_color = PALETTE["white"]
        else:
            add_rect(slide, x, y, w, h, PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        add_textbox(slide, x, y + Inches(0.03), w, h - Inches(0.06),
                    cell, font_size=10, color=text_color,
                    bold=(row_i == 0), alignment=PP_ALIGN.CENTER)

# Key findings
add_textbox(slide, Inches(7.0), Inches(4.2), Inches(5.7), Inches(0.4),
            "关键发现", font_size=14, color=PALETTE["dark_blue"], bold=True)

findings = [
    "● 0 m 和 +6 m 定位较好（误差 ≤1 m 超 74%）",
    "● -1 m 误差最大：仅 23 个样本，模型学习不充分",
    "● 两端距离向中间压缩（回归到均值现象）",
    "● 文件级 MAE > 片段级：简单平均不改善回归",
    "● R²=0.54：已学到距离信息，但精度仍有限",
]
add_multiline_textbox(slide, Inches(7.0), Inches(4.6), Inches(5.7), Inches(2.0),
                      findings, font_size=12, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

add_bottom_bar(slide)
add_page_number(slide, 15, TOTAL_SLIDES)

# ==============================================================
# Slide 16: Regression Baselines & Ablation
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "Stage1 基线对比与消融实验", "Baselines Comparison & Ablation Study")

# Baselines
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.4),
            "传统基线对比", font_size=16, color=PALETTE["dark_blue"], bold=True)

base_reg = [
    ["方法", "片段MAE", "片段RMSE", "≤1.0m%"],
    ["1D CNN", "1.90", "2.43", "35%"],
    ["SVM", "1.51 ★", "1.91 ★", "46%"],
    ["KNN", "1.76", "2.29", "44%"],
    ["决策树", "2.44", "2.97", "34%"],
]
for row_i, row in enumerate(base_reg):
    for col_j, cell in enumerate(row):
        x = Inches(0.7 + col_j * 1.45)
        y = Inches(2.0 + row_i * 0.48)
        w = Inches(1.4)
        h = Inches(0.43)
        if row_i == 0:
            add_rect(slide, x, y, w, h, PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        elif row_i == 2:
            add_rect(slide, x, y, w, h, PALETTE["accent2"])
            text_color = PALETTE["white"]
        else:
            add_rect(slide, x, y, w, h, PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        add_textbox(slide, x + Inches(0.05), y + Inches(0.05), w - Inches(0.1), h - Inches(0.1),
                    cell, font_size=13, color=text_color,
                    bold=(row_i == 0 or row_i == 2), alignment=PP_ALIGN.CENTER)

# Channel ablation
add_textbox(slide, Inches(0.7), Inches(4.1), Inches(5.8), Inches(0.4),
            "通道消融实验", font_size=16, color=PALETTE["dark_blue"], bold=True)

abl = [
    ["设置", "片段MAE", "片段RMSE", "≤1.0m%"],
    ["双通道（默认）", "1.90", "2.43", "35%"],
    ["单通道", "3.88 ✗", "4.75 ✗", "17%"],
]
for row_i, row in enumerate(abl):
    for col_j, cell in enumerate(row):
        x = Inches(0.7 + col_j * 1.45)
        y = Inches(4.6 + row_i * 0.48)
        w = Inches(1.4)
        h = Inches(0.43)
        if row_i == 0:
            add_rect(slide, x, y, w, h, PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        elif row_i == 2:
            add_rect(slide, x, y, w, h, RGBColor(0xE7, 0x4C, 0x3C))
            text_color = PALETTE["white"]
        else:
            add_rect(slide, x, y, w, h, PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        add_textbox(slide, x + Inches(0.05), y + Inches(0.05), w - Inches(0.1), h - Inches(0.1),
                    cell, font_size=13, color=text_color,
                    bold=(row_i >= 1), alignment=PP_ALIGN.CENTER)

# CNN vs Conformer for regression
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
            "CNN vs Conformer 回归对比", font_size=16, color=PALETTE["dark_blue"], bold=True)

cnn_conf_reg = [
    ["模型", "片段MAE", "片段RMSE", "文件MAE", "文件RMSE"],
    ["1D CNN", "1.90", "2.43", "2.46", "2.94"],
    ["Conformer", "1.95", "2.63", "2.51", "3.10"],
]
for row_i, row in enumerate(cnn_conf_reg):
    for col_j, cell in enumerate(row):
        x = Inches(7.0 + col_j * 1.4)
        y = Inches(2.0 + row_i * 0.48)
        w = Inches(1.35)
        h = Inches(0.43)
        if row_i == 0:
            add_rect(slide, x, y, w, h, PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        else:
            add_rect(slide, x, y, w, h, PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        add_textbox(slide, x + Inches(0.05), y + Inches(0.05), w - Inches(0.1), h - Inches(0.1),
                    cell, font_size=12, color=text_color,
                    bold=(row_i == 0), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(7.0), Inches(3.1), Inches(5.7), Inches(0.4),
            "→ Conformer 回归未超 CNN！", font_size=14, color=PALETTE["accent"], bold=True)

# Summary insights
add_textbox(slide, Inches(7.0), Inches(3.7), Inches(5.7), Inches(0.4),
            "核心结论", font_size=16, color=PALETTE["dark_blue"], bold=True)

summary_reg = [
    "① SVM 在当前回归任务上表现最强",
    "   → 手工时频特征 + 核方法仍具竞争力",
    "② 单通道 MAE 从 1.90 → 3.88 (+104%)",
    "   → 双通道差异是定位的关键信息来源",
    "③ Conformer 回归未优于 CNN",
    "   → 距离回归对细粒度差异敏感",
    "   → 强全局建模可能平滑局部关键差异",
    "④ 回归任务仍是后续优化的重点方向",
]
add_multiline_textbox(slide, Inches(7.0), Inches(4.1), Inches(5.7), Inches(2.8),
                      summary_reg, font_size=12, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

add_bottom_bar(slide)
add_page_number(slide, 16, TOTAL_SLIDES)

# ==============================================================
# Slide 17: Overall Results Summary
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "实验结果总览", "Overall Results Summary")

# Big summary table
summary_data = [
    ["任务", "模型", "片段主指标", "文件主指标", "状态"],
    ["Stage2\n分类", "1D CNN", "Acc 0.9415\nF1 0.9180", "Acc 0.9783\nF1 0.9836", "✓ 良好"],
    ["Stage2\n分类", "Conformer", "Acc 0.9900\nF1 0.9883", "Acc 1.0000\nF1 1.0000", "★ 最优"],
    ["Stage2\n分类", "SVM基线", "Acc 0.9300\nF1 0.9046", "Acc 1.0000\nF1 1.0000", "✓ 强基线"],
    ["Stage1\n回归", "1D CNN", "MAE 1.90\nRMSE 2.43", "MAE 2.46\nRMSE 2.94", "△ 可接受"],
    ["Stage1\n回归", "SVM基线", "MAE 1.51\nRMSE 1.91", "MAE 1.83\nRMSE 2.20", "★ 最优"],
    ["Stage1\n回归", "Conformer", "MAE 1.95\nRMSE 2.63", "MAE 2.51\nRMSE 3.10", "未超越CNN"],
]

for row_i, row in enumerate(summary_data):
    for col_j, cell in enumerate(row):
        col_widths = [Inches(1.5), Inches(1.8), Inches(3.0), Inches(3.0), Inches(2.3)]
        x = Inches(0.7 + sum(w.inches for w in col_widths[:col_j]))
        y = Inches(1.4 + row_i * 0.88)
        w = col_widths[col_j]
        h = Inches(0.8)

        if row_i == 0:
            add_rect(slide, x, y, w, h, PALETTE["dark_blue"])
            text_color = PALETTE["white"]
        elif "★" in cell:
            add_rect(slide, x, y, w, h, PALETTE["accent2"])
            text_color = PALETTE["white"]
        elif "✓" in cell:
            add_rect(slide, x, y, w, h, PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]
        elif "△" in cell:
            add_rect(slide, x, y, w, h, RGBColor(0xF3, 0x9C, 0x12))
            text_color = PALETTE["white"]
        else:
            add_rect(slide, x, y, w, h, PALETTE["light_gray"])
            text_color = PALETTE["dark_gray"]

        add_textbox(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), h - Inches(0.1),
                    cell, font_size=12, color=text_color,
                    bold=(row_i == 0), alignment=PP_ALIGN.CENTER)

# Bottom takeaway
add_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.1), PALETTE["dark_blue"])
takeaway = [
    "结论：① 分类任务已达到工程可用水平（文件级 Acc ≥ 97.8%） ② 距离回归验证了可行性但精度待提升（MAE 1.9 m）",
    "③ 双通道输入对定位至关重要（单通道 MAE 恶化 104%） ④ 更强模型 ≠ 更好回归（需匹配任务性质与数据规模）",
]
add_multiline_textbox(slide, Inches(0.9), Inches(6.05), Inches(11.5), Inches(1.0),
                      takeaway, font_size=14, color=PALETTE["white"],
                      line_spacing=1.7)

add_bottom_bar(slide)
add_page_number(slide, 17, TOTAL_SLIDES)

# ==============================================================
# Slide 18: Contributions & Limitations
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "主要贡献与研究局限", "Contributions & Limitations")

# Contributions
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.4),
            "主要贡献", font_size=18, color=PALETTE["accent2"], bold=True)

contribs = [
    "① 数据层面：构建了从原始 CSV → 分段样本的\n    自动化流程，文件级划分防止数据泄漏",
    "② 任务层面：将泄漏检测拆分为分类 + 回归\n    两阶段，降低耦合度，提高可解释性",
    "③ 模型层面：1D CNN 基础模型 + Conformer\n    结构对比，分析不同建模方式的适用边界",
    "④ 实验层面：多层级评测 + 传统基线 +\n    消融实验，避免单一指标误导结论",
]
add_multiline_textbox(slide, Inches(0.7), Inches(2.0), Inches(5.8), Inches(4.0),
                      contribs, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.8)

# Limitations
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
            "当前局限", font_size=18, color=RGBColor(0xE7, 0x4C, 0x3C), bold=True)

limits = [
    "① 数据规模有限：工况覆盖不足，距离档位\n"
    "    分布不均衡（-1 m 仅 23 样本）",
    "② 回归误差仍大：P95 误差 4.99 m，部分\n"
    "    距离容差命中率低于 10%",
    "③ 深度模型未全面领先：SVM 在回归任务\n"
    "    上 MAE 低于 1D CNN（1.51 vs 1.90）",
    "④ 文件级聚合策略简单：回归中简单平均\n"
    "    反而可能增大误差",
    "⑤ 未解决在线漂移：跨工况泛化、传感器\n"
    "    老化和环境变化有待验证",
]
add_multiline_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.7), Inches(5.0),
                      limits, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.8)

add_bottom_bar(slide)
add_page_number(slide, 18, TOTAL_SLIDES)

# ==============================================================
# Slide 19: Future Work & Conclusion
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])
add_section_header(slide, "结论与展望", "Conclusion & Future Work")

# Conclusion
add_rect(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.8), PALETTE["light_gray"])
add_textbox(slide, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.4),
            "全文总结", font_size=18, color=PALETTE["dark_blue"], bold=True)

conclusions = [
    "1. 验证了两阶段框架的可行性：Stage2 分类 + Stage1 回归，任务拆解清晰，模块可独立优化",
    "2. 分类任务表现成熟：CNN 文件级 Acc 97.8%，Conformer 达到文件级 100%，具备工程可用性",
    "3. 距离回归初步验证：CNN 片段级 MAE 1.90 m，证明双通道原始波形中包含可学习的距离信息",
    "4. 双通道输入是定位基础：消融实验证实单通道退化严重（MAE +104%），通道间差异 = 定位关键",
    "5. 深度模型 ≠ 全面更优：SVM 在回归任务上领先，提醒模型选择需匹配任务性质与数据规模",
]
add_multiline_textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.2),
                      conclusions, font_size=13, color=PALETTE["dark_gray"],
                      line_spacing=1.7)

# Future work
add_textbox(slide, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.4),
            "后续工作展望", font_size=18, color=PALETTE["accent"], bold=True)

future = [
    ("数据扩充", "补充多工况、多距离档位数据\n跨压力/流量/管径实验验证"),

    ("模型改进", "距离感知损失、时频特征融合\n不确定性估计、轻量化部署"),

    ("评测升级", "跨批次/跨工况严格测试\n在线增量学习与漂移检测"),
]

for i, (title, desc) in enumerate(future):
    x = Inches(0.7 + i * 4.1)
    add_rect(slide, x, Inches(5.1), Inches(3.8), Inches(1.5), PALETTE["light_gray"])
    add_textbox(slide, x + Inches(0.15), Inches(5.15), Inches(3.5), Inches(0.35),
                title, font_size=15, color=PALETTE["accent"], bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(5.55), Inches(3.5), Inches(0.9),
                desc, font_size=12, color=PALETTE["dark_gray"])

add_bottom_bar(slide)
add_page_number(slide, 19, TOTAL_SLIDES)

# ==============================================================
# Slide 20: Thank You
# ==============================================================
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, PALETTE["white"])

# Large dark blue block
add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PALETTE["dark_blue"])
add_rect(slide, Inches(0), Inches(3.3), SLIDE_W, Inches(0.08), PALETTE["accent"])

add_textbox(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(1.2),
            "感谢各位老师批评指正",
            font_size=40, color=PALETTE["white"], bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(3.8), SLIDE_W, Inches(0.8),
            "基于双通道时序信号的两阶段管道泄漏检测与定位方法",
            font_size=18, color=PALETTE["gray"],
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(4.6), SLIDE_W, Inches(0.5),
            "答辩人：孟达  |  天津大学建筑工程学院  |  船舶与海洋工程",
            font_size=14, color=PALETTE["gray"],
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.5),
            "导师：顾鹏 讲师",
            font_size=14, color=PALETTE["gray"],
            alignment=PP_ALIGN.CENTER)

# Q&A note
add_textbox(slide, Inches(0), Inches(6.2), SLIDE_W, Inches(0.5),
            "欢迎提问  ·  Q & A",
            font_size=22, color=PALETTE["accent"], bold=True,
            alignment=PP_ALIGN.CENTER)

# ==============================================================
# Save
# ==============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "defense_presentation.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
