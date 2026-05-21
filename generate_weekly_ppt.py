"""Generate a placeholder weekly report deck for refreshed experiments."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("outputs/本周工作汇报PPT_placeholder.pptx")


def add_slide(prs: Presentation, title: str, body: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(8.5), Inches(3.4))
    frame = box.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = frame.add_paragraph()
    p2.text = body
    p2.font.size = Pt(17)
    p2.alignment = PP_ALIGN.CENTER


def main() -> None:
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    add_slide(
        prs,
        "本周工作汇报占位",
        "旧数据集实验结果已清理，等待 Internoise 新训练结果。",
    )
    add_slide(
        prs,
        "后续更新项",
        "Stage2 五分类指标、Stage1 回归误差、训练曲线、混淆矩阵和论文表格。",
    )
    prs.save(OUTPUT_PATH)
    print(f"Placeholder weekly PPT written to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
